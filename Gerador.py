import sys
import subprocess
import importlib

# Dicionário: "nome para importar" : "nome do pacote no PyPI"
dependencias = {
    "ttkbootstrap": "ttkbootstrap",
    "pptx": "python-pptx",
    "googleapiclient": "google-api-python-client",
    "google.auth.transport": "google-auth-httplib2",
    "google_auth_oauthlib": "google-auth-oauthlib",
    "requests": "requests"
}

for modulo, pacote in dependencias.items():
    try:
        importlib.import_module(modulo)
    except ImportError:
        print(f"Instalando '{pacote}' pois o módulo '{modulo}' não foi encontrado...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", pacote])
        # Atualiza o cache de importação
        importlib.invalidate_caches()
        try:
            importlib.import_module(modulo)
        except ImportError:
            print(f"Falha ao importar '{modulo}' mesmo após a instalação.")



import os
import requests
import io
import pickle
import json
from datetime import date

# Tkinter e ttkbootstrap
import tkinter as tk
from tkinter import ttk
import ttkbootstrap as ttkb
from tkinter.messagebox import showerror, showinfo

# python-pptx
from pptx import Presentation

# Google Drive / Auth
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload





def baixar_arquivo_if_needed(nome_arquivo, url):
    if not os.path.exists(nome_arquivo):
        print(f"Baixando {nome_arquivo}...")
        r = requests.get(url)
        with open(nome_arquivo, "wb") as f:
            f.write(r.content)
            

# ---------------------------------------------------------
# Ajustes Globais
# ---------------------------------------------------------
script_dir = os.path.dirname(os.path.abspath(__file__))
os.chdir(script_dir)
print("Novo diretório de trabalho:", os.getcwd())

CONFIG_FILE = "config_vendedor.json"
MAX_ABAS = 10


# ---------------------------------------------------------
# Configurações de vendedor (salvar/carregar)
# ---------------------------------------------------------
def carregar_config(nome_closer_var, celular_closer_var, email_closer_var):
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            nome_closer_var.set(data.get("nome_vendedor", ""))
            celular_closer_var.set(data.get("celular_vendedor", ""))
            email_closer_var.set(data.get("email_vendedor", ""))
        except json.JSONDecodeError:
            pass

def salvar_config(nome_closer, celular_closer, email_closer):
    dados = {
        "nome_vendedor": nome_closer,
        "celular_vendedor": celular_closer,
        "email_vendedor": email_closer
    }
    if os.path.exists(CONFIG_FILE):
        os.chmod(CONFIG_FILE, 0o666)
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(dados, f, indent=4, ensure_ascii=False)
    except PermissionError:
        pass


# ---------------------------------------------------------
# Dados de Planos e Tabelas de Preço
# ---------------------------------------------------------

PLAN_INFO = {
    "Plano PDV": {
        "base_mensal": 110.00,
        "base_users": 2,
        "max_additional_users": 1,
        "user_extra_cost": 19.00,
        "base_pdv": 1,
        "max_additional_pdv": 0,
        "pdv_extra_cost": 0.00,
        "mandatory": [
            "30 Notas Fiscais", # Este é parte do plano, não um checkbox opcional usual
            "Suporte Técnico - Via Chamados",
            "Relatório Básico",
        ],
        "allowed_optionals": [
            "Smart Menu", 
            "Delivery",
            "Hub de Delivery",
            "Delivery Direto Profissional",
            "Delivery Direto VIP", 
            "TEF",
            "Importação de XML",
            "Cardápio digital",
            "Autoatendimento",
            "Contratos de cartões e outros", 
            "Ordem de Serviço", 
            "Estoque em Grade",
            "Conciliação Bancária"

        ],
        "module_limits": {}
    },
    "Gestão": {
        "base_mensal": 221.10,
        "base_users": 3,
        "max_additional_users": 2,
        "user_extra_cost": 19.00,
        "base_pdv": 2,
        "max_additional_pdv": 1,
        "pdv_extra_cost": 59.90,
        "mandatory": [
            "Notas Fiscais Ilimitadas",
            "Importação de XML",
            "Estoque em Grade",
            "Relatório básico",
            "Suporte Técnico - Via Chamados",
            "Suporte Técnico - Via Chat",
            "Delivery"
        ],
        "allowed_optionals": [
            "Relatório KDS",
            "Integração API", # Permitido para o Plano Gestão
            "Suporte Técnico - Estendido", 
            "Conciliação Bancária", 
            "Contratos de cartões e outros",
            "Controle de Mesas", 
            "Cardápio digital",
            "Ordem de Serviço", 
            "Produção", 
            "Relatório Dinâmico",
            "Backup Realtime", 
            "Business Intelligence (BI)", 
            "Hub de Delivery", "Facilita NFE",
            "App Gestão CPlug",
            "Smart Menu", 
            "API DD", #API Delivery Direto, não confundir com Integração API, adicionar depois
            "Delivery Direto Básico", 
            "Delivery Direto Profissional", 
            "Delivery Direto VIP",
            "Promoções", 
            "Marketing", 
            "Painel de Senha Mobile", 
            "Painel de Senha TV",
            "Domínio Próprio",
            "Entrega Fácil iFood", #ajustar preço
            "Gestão de Redes Sociais", 
            "Combo de Logística", 
            "Painel MultiLojas",
            "Central Telefônica",
            "E-mail Profissional",
            "Atualização em Tempo Real"
        ],
        "module_limits": {
            }, # Limites de módulos específicos, se necessário. por exemplo, "Smart TEF": 5
        "fixed_quantities": {}
    },
    "Performance": {
        "base_mensal": 554.50,
        "base_users": 5,
        "max_additional_users": 5,
        "user_extra_cost": 19.00,
        "base_pdv": 3,
        "max_additional_pdv": 2,
        "pdv_extra_cost": 59.90,
        "Backup Realtime": 99.00,
        "mandatory": [
            "Produção",
            "Promoções",
            "Notas Fiscais Ilimitadas",
            "Importação de XML",
            "Hub de Delivery",
            "Ordem de Serviço",
            "Delivery",
            "App Gestão CPlug", # App Gestão CPlug é spinbox
            "Relatório KDS",
            "Painel de Senha TV",
            "Painel de Senha Mobile",
            "Controle de Mesas",
            "Estoque em Grade",
            "Marketing",
            "Relatório Básico",
            "Relatório Dinâmico",
            "Atualização em Tempo Real",
            "Facilita NFE",
            "Conciliação Bancária",
            "Contratos de cartões e outros",
            "Suporte Técnico - Via Chamados",
            "Suporte Técnico - Via Chat",
            "Suporte Técnico - Estendido"
        ],
        "allowed_optionals": [
            "Integração API",
            "Integração TAP",
            "Backup Realtime",
            "Business Intelligence (BI)",
            "Smart Menu",
            "Programa de Fidelidade",
            "Delivery Direto Profissional",
            "Delivery Direto VIP",
            "Domínio Próprio",
            "Entrega Fácil iFood",
            "Robô de WhatsApp + Recuperador de Pedido",
            "Gestão de Redes Sociais",
            "Combo de Logística",
            "Painel MultiLojas",
            "Central Telefônica",
            "API DD",
            "E-mail Profissional",
            "Cardápio digital",
            "Delivery Direto Básico",
            # Adicione outros módulos de checkbox que o Plano Performance pode ter como opcionais (se houver algum não mandatório)
        ],
        "module_limits": {}, # Limites de módulos específicos, se necessário. por exemplo, "Smart TEF": 5;
        "fixed_quantities": {
            #AQUI TU PODE DEFINIR QUANTIDADE FIXA DE USUÁRIOS, PDVS, ENFIM, OS MÓDULOS SE HOUVER
        }
    },
    # ... Defina "allowed_optionals": [] para os outros planos ("Autoatendimento", "Bling", "Em Branco")
    #     e popule essas listas conforme as regras de negócio para cada um.
    #     Por exemplo, para "Em Branco", talvez todos os opcionais sejam permitidos:
    "Em Branco": {
        "base_mensal": 0.00, # Exemplo
        "base_users": 1,
        "max_additional_users": 10,
        "user_extra_cost": 19.00,
        "base_pdv": 0,
        "max_additional_pdv": 10,
        "pdv_extra_cost": 59.90,
        "mandatory": [],
        "allowed_optionals": [ # Lista todos os módulos de self.modules que podem ser adicionados
            "Suporte Técnico - Via Chamados",
            "Suporte Técnico - Via Chat",
            "Suporte Técnico - Estendido",
            "Conciliação Bancária",
            "Contratos de cartões e outros",
            "Controle de Mesas",
            "Delivery",
            "Estoque em Grade",
            "Importação de XML",
            "Ordem de Serviço",
            "Produção",
            "Relatório Dinâmico",
            "Notas Fiscais Ilimitadas",
            "Backup Realtime",
            "Atualização em Tempo Real",
            "Business Intelligence (BI)",
            "App Gestão CPlug",
            "Hub de Delivery",
            "Facilita NFE",
            "Smart Menu",
            "Programa de Fidelidade",
            "Delivery Direto Profissional",
            "Delivery Direto VIP",
            "Promoções",
            "Marketing",
            "Painel de Senha TV",
            "Painel de Senha Mobile",
            "Relatório KDS",
            "Integração TAP",
            "Integração API",
            "Entrega Fácil iFood",
            "Painel MultiLojas",
            "Central Telefônica",
            "Relatório Básico",
            "30 Notas Fiscais" # Mesmo que sejam mandatórios em outros planos
        ],
        "module_limits": {}
    },
    # Adicione allowed_optionals para "Autoatendimento" e "Bling"
    "Autoatendimento": { # Exemplo, ajuste conforme necessário
        "base_mensal": 0.0, # Verifique a precificação correta
        # ... outras chaves ...
        "mandatory": ["Suporte Técnico - Via Chamados"], # Exemplo
        "allowed_optionals": [
            "Suporte Técnico - Estendido", "Smart Menu", "Integração API" # etc.
        ],
        "module_limits": {}
    },
    "Bling": { # Exemplo, ajuste conforme necessário
        "base_mensal": 229.90,
        "base_anual": 179.90,
        "base_users": 2,
        "max_additional_users": 5,
        "user_extra_cost": 19.00,
        "base_pdv": 1,
        "max_additional_pdv": 5,
        "pdv_extra_cost": 50.00,
        "mandatory": [
            "Suporte Técnico - Via Chamados",
            "Relatório Básico",
            "Notas Fiscais Ilimitadas",
            ], # Exemplo
        "allowed_optionals": [
            "Suporte Técnico - Estendido", 
            "Delivery",
            "Contratos de cartões e outros",
            "Estoque em Grade",
            "Controle de Mesas" # etc.
        ],
        "module_limits": {}
    }
}



SEM_DESCONTO = {
    "TEF",
    "Autoatendimento",
    "Smart TEF",
    "Domínio Próprio",
    "Entrega Fácil iFood",
    "Robô de WhatsApp + Recuperador de Pedido",
    "Gestão de Redes Sociais",
    "Combo de Logística",
    "Painel MultiLojas","Programa de Fidelidade","Integração API", "Integração TAP",
    "Central Telefônica"
}

precos_mensais = {
    "30 Notas Fiscais": 0.00, # Mantido, pois é parte de um plano fixo e não opcional com novo preço
    "Conciliação Bancária": 50.00,
    "Contratos de cartões e outros": 50.00,
    "Controle de Mesas": 49.00,
    "Delivery": 30.00, # Mantido, pois é parte de planos fixos e não opcional com novo preço individual listado
    "Estoque em Grade": 40.00,
    "Importação de XML": 29.00,
    "Ordem de Serviço": 20.00, # O último preço encontrado foi R$20,00 (Plano Gestão)
    "Produção": 30.00,
    "Relatório Dinâmico": 50.00,
    "Notas Fiscais Ilimitadas": 119.90, # Mantido, pois é parte de planos fixos e não opcional com novo preço individual listado

    "60 Notas Fiscais": 40.00, # Mantido da lista original
    "120 Notas Fiscais": 70.00, # Mantido da lista original
    "250 Notas Fiscais": 90.00, # Mantido da lista original

    "TEF": 99.00, # Encontrado como 99,90 e 99,00. Usado 99,90 (mais frequente/último no Plano Performance)
    "Smart TEF": 49.90,
    "Backup Realtime": 199.00, # O último preço encontrado foi R$99,90 (Plano Performance)
    "Atualização em Tempo Real": 49.00, # Nome e preço atualizados
    "Business Intelligence (BI)": 99.00, # O último preço encontrado foi R$99,00 (Plano Performance)
    "Hub de Delivery": 79.00,
    "Facilita NFE": 99.00,
    "Smart Menu": 99.00, # O último preço encontrado foi R$99,90 (Plano Gestão/Performance)
    "Cardápio digital": 99.00, # Nome e preço atualizados
    "Programa de Fidelidade": 299.90,
    "API DD": 49.90, # API Delivery Direto, não confundir com Integração API
    "Autoatendimento": 299.00,
    "Delivery Direto Básico": 99.00, # Mantido da lista original
    "Delivery Direto Profissional": 200.00,
    "Delivery Direto VIP": 300.00,
    "Promoções": 24.50,
    "Marketing": 24.50,
    "Painel de Senha TV": 00.00, # Mantido da lista original (Painel Senha TV/Mobile são itens separados na nova lista)
    "Integração TAP": 299.00, # Nome e preço atualizados
    "Integração API": 299.00,
    "Relatório KDS": 29.90, # Mantido, pois é parte de planos fixos e não opcional com novo preço individual listado
    "App Gestão CPlug": 20.00,

    "Domínio Próprio": 19.90, # Mantido da lista original
    "Entrega Fácil iFood": 19.90, # Mantido da lista original
    "Painel MultiLojas": 199.00, # Mantido da lista original
    "Central Telefônica": 399.90, # Mantido da lista original
    # Novos itens adicionados da lista fornecida
    "Entrega Fácil iFood": 49.90,
    "E-mail Profissional": 19.90,
    

    "Painel de Senha Mobile": 49.00,
    "Painel de Senhta TV": 00.00,
    "Suporte Técnico - Estendido": 99.00,
    "Suporte Técnico - Via Chat": 00.00
    
}

CUSTOS_SETUP_ADICIONAL = {
    "Programa de Fidelidade": 1000.00,
    "Central Telefônica": 300.00
}

# ---------------------------------------------------------
# Função utilitária para substituir placeholders no Slide
# ---------------------------------------------------------
def substituir_placeholders_no_slide(slide, dados):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    txt = run.text
                    for k, v in dados.items():
                        if k in txt:
                            txt = txt.replace(k, v)
                    run.text = txt


# ---------------------------------------------------------
# Classe PlanoFrame (Aba)
#   => contém toda a lógica de cálculo do plano
# ---------------------------------------------------------
class PlanoFrame(ttkb.Frame):
    def __init__(
        self,
        parent,
        aba_index,
        nome_cliente_var_shared,
        validade_proposta_var_shared,
        on_close_callback=None
    ):
        super().__init__(parent)
        self.aba_index = aba_index
        self.on_close_callback = on_close_callback

        # Variáveis compartilhadas para Nome do Cliente, Validade e Plano
        self.nome_cliente_var = nome_cliente_var_shared
        self.validade_proposta_var = validade_proposta_var_shared
        self.nome_plano_var = tk.StringVar(value="") # ← Nome do plano

        self.current_plan = "Plano PDV"
        self.spin_pdv_var = tk.IntVar(value=1)
        self.spin_users_var = tk.IntVar(value=1)
        self.spin_auto_var = tk.IntVar(value=0)
        self.spin_cardapio_var = tk.IntVar(value=0)
        self.spin_tef_var = tk.IntVar(value=0)
        self.spin_smart_tef_var = tk.IntVar(value=0)
        self.spin_app_cplug_var = tk.IntVar(value=0)
        self.spin_delivery_direto_basico_var = tk.IntVar(value=0)
        self.var_notas = tk.StringVar(value="NONE")

        # Módulos extras (checkboxes)
        self.modules = {}
        self.check_buttons = {}

        # Overrides de cálculo
        self.user_override_anual_active = tk.BooleanVar(value=False)
        self.user_override_discount_active = tk.BooleanVar(value=False)
        self.valor_anual_editavel = tk.StringVar(value="0.00")
        self.desconto_personalizado = tk.StringVar(value="0")


        # Armazenar valores
        self.computed_mensal = 0.0
        self.computed_anual = 0.0
        self.computed_desconto_percent = 0.0

        # Layout com scrollbar
        self.canvas = tk.Canvas(self)
        self.canvas.pack(side="left", fill="both", expand=True)
        self.scrollbar = ttkb.Scrollbar(self, orient="vertical", command=self.canvas.yview)
        self.scrollbar.pack(side="right", fill="y")
        self.canvas.configure(yscrollcommand=self.scrollbar.set)
        self.container = ttkb.Frame(self.canvas)
        self.canvas.create_window((0,0), window=self.container, anchor="nw")
        self.container.bind("<Configure>", lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all")))

        self.frame_main = ttkb.Frame(self.container)
        self.frame_main.pack(fill="both", expand=True)

        self.frame_left = ttkb.Frame(self.frame_main)
        self.frame_left.pack(side="left", fill="both", expand=True, padx=5, pady=5)

        self.frame_right = ttkb.Frame(self.frame_main)
        self.frame_right.pack(side="left", fill="y", padx=5, pady=5)

        self._montar_layout_esquerda()
        self._montar_layout_direita()
        self.configurar_plano("Plano PDV")

    def fechar_aba(self):
        if self.on_close_callback:
            self.on_close_callback(self.aba_index)

    def _montar_layout_esquerda(self):
        top_bar = ttkb.Frame(self.frame_left)
        top_bar.pack(fill="x", pady=5)
        ttkb.Label(top_bar, text=f"Aba Plano {self.aba_index}", font="-size 12 -weight bold").pack(side="left")
        btn_close = ttkb.Button(top_bar, text="Fechar Aba", command=self.fechar_aba)
        btn_close.pack(side="right")

        frame_planos = ttkb.Labelframe(self.frame_left, text="Planos")
        frame_planos.pack(fill="x", pady=5)
        for p in ["Plano PDV","Gestão","Performance","Autoatendimento","Bling", "Em Branco"]:
            ttkb.Button(frame_planos, text=p,
                        command=lambda pl=p: self.configurar_plano(pl)
                       ).pack(side="left", padx=5)

        frame_notas = ttkb.Labelframe(self.frame_left, text="Notas Fiscais")
        frame_notas.pack(fill="x", pady=5)
        f_nf = ttkb.Frame(frame_notas)
        f_nf.pack(fill="x", padx=5, pady=5)
        for nfopt in ["60","120","250"]:
            rb = ttk.Radiobutton(f_nf, text=nfopt+" Notas",
                                 variable=self.var_notas, value=nfopt,
                                 command=self.atualizar_valores)
            rb.pack(side="left", padx=5)

        # Módulos (checkboxes)
        self.modules = {
            "30 Notas Fiscais": tk.IntVar(),
            "Relatório Básico": tk.IntVar(),
            "Suporte Técnico - Via Chamados": tk.IntVar(),
            "Suporte Técnico - Via Chat": tk.IntVar(),
            "Suporte Técnico - Estendido": tk.IntVar(),
            "Conciliação Bancária": tk.IntVar(),
            "Contratos de cartões e outros": tk.IntVar(),
            "Cardápio digital": tk.IntVar(),
            "Controle de Mesas": tk.IntVar(),
            "Delivery": tk.IntVar(),
            "Estoque em Grade": tk.IntVar(),
            "Importação de XML": tk.IntVar(),
            "Ordem de Serviço": tk.IntVar(),
            "Produção": tk.IntVar(),
            "Relatório Dinâmico": tk.IntVar(),
            "Notas Fiscais Ilimitadas": tk.IntVar(),
            "Backup Realtime": tk.IntVar(),
            "Atualização em Tempo Real": tk.IntVar(),
            "Business Intelligence (BI)": tk.IntVar(),
            "Hub de Delivery": tk.IntVar(),
            "Facilita NFE": tk.IntVar(),
            "Smart Menu": tk.IntVar(),
            "Programa de Fidelidade": tk.IntVar(),
            "Delivery Direto Básico": tk.IntVar(),
            "API DD": tk.IntVar(),  # API Delivery Direto, não confundir com Integração API
            "Delivery Direto Profissional": tk.IntVar(),
            "Delivery Direto VIP": tk.IntVar(),
            "Promoções": tk.IntVar(),
            "Marketing": tk.IntVar(),
            "Painel de Senha TV": tk.IntVar(),
            "Painel de Senha Mobile": tk.IntVar(),
            "Relatório KDS": tk.IntVar(),
            "Integração TAP": tk.IntVar(),
            "Integração API": tk.IntVar(),
            "Entrega Fácil iFood": tk.IntVar(),
            "E-mail Profissional": tk.IntVar(),
            "Painel MultiLojas": tk.IntVar(),
            "Central Telefônica": tk.IntVar(),
            "App Gestão CPlug": tk.IntVar(),
            "Ordem de Serviço": tk.IntVar(),
        }

        frame_mod = ttkb.Labelframe(self.frame_left, text="Outros Módulos")
        frame_mod.pack(fill="both", expand=True, pady=5)
        f_mod_cols = ttkb.Frame(frame_mod)
        f_mod_cols.pack(fill="both", expand=True)

        f_mod_left = ttkb.Frame(f_mod_cols)
        f_mod_left.pack(side="left", fill="both", expand=True, padx=5)
        f_mod_right = ttkb.Frame(f_mod_cols)
        f_mod_right.pack(side="left", fill="both", expand=True, padx=5)

        all_mods = sorted(self.modules.keys())
        mid = len(all_mods)//2
        left_side = all_mods[:mid]
        right_side = all_mods[mid:]
        self.check_buttons = {}

        for m in left_side:
            cb = ttk.Checkbutton(f_mod_left, text=m,
                                 variable=self.modules[m],
                                 command=self.atualizar_valores)
            cb.pack(anchor="w", pady=2)
            self.check_buttons[m] = cb

        for m in right_side:
            cb = ttk.Checkbutton(f_mod_right, text=m,
                                 variable=self.modules[m],
                                 command=self.atualizar_valores)
            cb.pack(anchor="w", pady=2)
            self.check_buttons[m] = cb

        frame_dados = ttkb.Labelframe(self.frame_left, text="Dados do Cliente")
        frame_dados.pack(fill="x", pady=5)
        ttkb.Label(frame_dados, text="Nome do Cliente:").grid(row=0, column=0, sticky="w")
        ttkb.Entry(frame_dados, textvariable=self.nome_cliente_var).grid(row=0, column=1, padx=5, pady=2)
        ttkb.Label(frame_dados, text="Validade Proposta:").grid(row=1, column=0, sticky="w")
        ttkb.Entry(frame_dados, textvariable=self.validade_proposta_var).grid(row=1, column=1, padx=5, pady=2)

        ttkb.Label(frame_dados, text="Nome do Plano:").grid(row=2, column=0, sticky="w")
        ttkb.Entry(frame_dados, textvariable=self.nome_plano_var).grid(row=2, column=1, padx=5, pady=2)



    def _montar_layout_direita(self):
        frame_inc = ttkb.Labelframe(self.frame_right, text="Incrementos")
        frame_inc.pack(fill="x", pady=5)

        ttkb.Label(frame_inc, text="PDVs").grid(row=0, column=0, sticky="w")
        self.sp_pdv = ttkb.Spinbox(frame_inc, from_=0, to=99,
                              textvariable=self.spin_pdv_var,
                              command=self.atualizar_valores)
        self.sp_pdv.grid(row=0, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="Usuários").grid(row=1, column=0, sticky="w")
        self.sp_usr = ttkb.Spinbox(frame_inc, from_=0, to=999,
                              textvariable=self.spin_users_var,
                              command=self.atualizar_valores)
        self.sp_usr.grid(row=1, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="Autoatendimento").grid(row=2, column=0, sticky="w")
        self.sp_at = ttkb.Spinbox(frame_inc, from_=0, to=999,
                             textvariable=self.spin_auto_var,
                             command=self.atualizar_valores)
        self.sp_at.grid(row=2, column=1, padx=5, pady=2)

# AQUI ERA O CARDAPIO DIGITAL NO SPIN BOX., PARA VOLTAR, É SÓ TIRAR DO COMENTARIO.
#        ttkb.Label(frame_inc, text="Cardápio Digital").grid(row=3, column=0, sticky="w")
#        self.sp_cd = ttkb.Spinbox(frame_inc, from_=0, to=999,
#                             textvariable=self.spin_cardapio_var,
#                             command=self.atualizar_valores)
#        self.sp_cd.grid(row=3, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="TEF").grid(row=4, column=0, sticky="w")
        self.sp_tef = ttkb.Spinbox(frame_inc, from_=0, to=99,
                              textvariable=self.spin_tef_var,
                              command=self.atualizar_valores)
        self.sp_tef.grid(row=4, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="Smart TEF").grid(row=5, column=0, sticky="w")
        self.sp_smf = ttkb.Spinbox(frame_inc, from_=0, to=99,
                              textvariable=self.spin_smart_tef_var,
                              command=self.atualizar_valores)
        self.sp_smf.grid(row=5, column=1, padx=5, pady=2)

        #ttkb.Label(frame_inc, text="App Gestão CPlug").grid(row=6, column=0, sticky="w")
        #self.sp_app = ttkb.Spinbox(frame_inc, from_=0, to=999,
        #                      textvariable=self.spin_app_cplug_var,
        #                      command=self.atualizar_valores)
        #self.sp_app.grid(row=6, column=1, padx=5, pady=2)

        #ttkb.Label(frame_inc, text="Delivery Direto Básico").grid(row=7, column=0, sticky="w")
        #self.sp_ddb = ttkb.Spinbox(frame_inc, from_=0, to=999,
        #                      textvariable=self.spin_delivery_direto_basico_var,
        #                      command=self.atualizar_valores)
        #self.sp_ddb.grid(row=7, column=1, padx=5, pady=2)

        frame_valores = ttkb.Labelframe(self.frame_right, text="Valores Finais")
        frame_valores.pack(fill="x", pady=5)

        self.lbl_plano_mensal = ttkb.Label(frame_valores, text="Plano (Mensal): R$ 0.00", font="-size 12 -weight bold")
        self.lbl_plano_mensal.pack()
        self.lbl_plano_anual = ttkb.Label(frame_valores, text="Plano (Anual): R$ 0.00", font="-size 12 -weight bold")
        self.lbl_plano_anual.pack()
        self.lbl_treinamento = ttkb.Label(frame_valores, text="Custo Treinamento (Mensal): R$ 0.00", font="-size 12 -weight bold")
        self.lbl_treinamento.pack()
        self.lbl_desconto = ttkb.Label(frame_valores, text="Desconto: 0%", font="-size 12 -weight bold")
        self.lbl_desconto.pack()
        
        self.lbl_custo_setup_adicional = ttkb.Label(frame_valores, text="Custo Adicional Setup: R$ 0,00", font="-size 12 -weight bold") # Pode ajustar o tamanho da fonte se desejar
        self.lbl_custo_setup_adicional.pack(pady=(1,0)) # Adiciona um pouco de espaço acima

        frame_edit_anual = ttkb.Labelframe(self.frame_right, text="Plano (Anual) (editável)")
        frame_edit_anual.pack(pady=5, fill="x")
        e_anual = ttkb.Entry(frame_edit_anual, textvariable=self.valor_anual_editavel, width=10)
        e_anual.pack(side="left", padx=5)
        e_anual.bind("<KeyRelease>", self.on_user_edit_valor_anual)
        b_reset_anual = ttkb.Button(frame_edit_anual, text="Reset Anual", command=self.on_reset_anual)
        b_reset_anual.pack(side="left", padx=5)

        frame_edit_desc = ttkb.Labelframe(self.frame_right, text="Desconto (%) (editável)")
        frame_edit_desc.pack(pady=5, fill="x")
        e_desc = ttkb.Entry(frame_edit_desc, textvariable=self.desconto_personalizado, width=10)
        e_desc.pack(side="left", padx=5)
        e_desc.bind("<KeyRelease>", self.on_user_edit_desconto)
        b_reset_desc = ttkb.Button(frame_edit_desc, text="Reset Desconto", command=self.on_reset_desconto)
        b_reset_desc.pack(side="left", padx=5)

    def on_user_edit_valor_anual(self, *args):
        self.user_override_anual_active.set(True)
        self.user_override_discount_active.set(False)
        self.atualizar_valores()

    def on_reset_anual(self):
        self.user_override_anual_active.set(False)
        self.valor_anual_editavel.set("0.00")
        self.atualizar_valores()

    def on_user_edit_desconto(self, *args):
        self.user_override_discount_active.set(True)
        self.user_override_anual_active.set(False)
        self.atualizar_valores()

    def on_reset_desconto(self):
        self.user_override_discount_active.set(False)
        self.desconto_personalizado.set("0")
        self.atualizar_valores()

    def configurar_plano(self, plano_selecionado):
            # Certifique-se que o plano_selecionado existe em PLAN_INFO
            if plano_selecionado not in PLAN_INFO:
                showerror("Erro de Configuração", f"Plano '{plano_selecionado}' não encontrado em PLAN_INFO.")
                return

            info = PLAN_INFO[plano_selecionado]
            self.current_plan = plano_selecionado
            self.nome_plano_var.set(plano_selecionado) # Atualiza o nome do plano na UI

            # --- Configuração de PDVs ---
            base_pdvs = info.get("base_pdv", 0)
            max_additional_pdvs = info.get("max_additional_pdv", 0)
            total_max_pdvs = base_pdvs + max_additional_pdvs

            self.spin_pdv_var.set(base_pdvs)
            if hasattr(self, 'sp_pdv'): # Verifica se a referência ao spinbox existe
                self.sp_pdv.config(from_=base_pdvs, to=total_max_pdvs)
                if max_additional_pdvs == 0:
                    self.sp_pdv.config(state='readonly' if base_pdvs > 0 else 'disabled')
                else:
                    self.sp_pdv.config(state='normal')
            else:
                print(f"Aviso: Referência self.sp_pdv não encontrada para o plano {plano_selecionado}.")

            # --- Configuração de Usuários ---
            base_users = info.get("base_users", 0)
            max_additional_users = info.get("max_additional_users", 0)
            total_max_users = base_users + max_additional_users

            self.spin_users_var.set(base_users)
            if hasattr(self, 'sp_usr'):
                self.sp_usr.config(from_=base_users, to=total_max_users)
                if max_additional_users == 0:
                    self.sp_usr.config(state='readonly' if base_users > 0 else 'disabled')
                else:
                    self.sp_usr.config(state='normal')
            else:
                print(f"Aviso: Referência self.sp_usr não encontrada para o plano {plano_selecionado}.")

            # --- Configuração de Módulos com Quantidade Controlada por Spinbox ---
            fixed_quantities = info.get("fixed_quantities", {})
            module_limits = info.get("module_limits", {})

            # Exemplo para Smart TEF (self.sp_smf)
            if hasattr(self, 'sp_smf'):
                if "Smart TEF" in fixed_quantities:
                    qty = fixed_quantities["Smart TEF"]
                    self.spin_smart_tef_var.set(qty)
                    self.sp_smf.config(from_=qty, to=qty, state='readonly')
                elif "Smart TEF" in module_limits:
                    limit = module_limits["Smart TEF"]
                    self.spin_smart_tef_var.set(0) # Opcional começa com 0
                    self.sp_smf.config(from_=0, to=limit, state='normal')
                else: # Sem limite específico ou fixo para Smart TEF neste plano
                    self.spin_smart_tef_var.set(0)
                    self.sp_smf.config(from_=0, to=99, state='normal') # Limite genérico
            else:
                # Se o spinbox não existir mas o módulo for relevante, pode haver um aviso
                if "Smart TEF" in fixed_quantities or "Smart TEF" in module_limits:
                    print(f"Aviso: Referência self.sp_smf não encontrada, mas Smart TEF é configurável para {plano_selecionado}.")


            # Adicione aqui lógica similar para outros módulos controlados por Spinbox
            # Ex: Autoatendimento (self.sp_auto), Cardápio Digital (self.sp_cardapio), etc.
            # if hasattr(self, 'sp_auto'):
            #     if "Autoatendimento" in fixed_quantities: ...
            #     elif "Autoatendimento" in module_limits: ...
            #     else: ...
            # Resetar spinboxes de módulos opcionais para 0 se não forem fixos ou limitados especificamente.
            if plano_selecionado == "Autoatendimento": # Lógica específica do seu código original
                self.spin_auto_var.set(1) # Ou o valor base do plano Autoatendimento
                # Configure os limites do self.sp_auto aqui se necessário
            elif "Autoatendimento" not in fixed_quantities and "Autoatendimento" not in module_limits:
                self.spin_auto_var.set(0)
                if hasattr(self, 'sp_auto'): self.sp_auto.config(from_=0, to=99, state='normal') # Limite genérico


            if "Cardápio Digital" not in fixed_quantities and "Cardápio Digital" not in module_limits:
                self.spin_cardapio_var.set(0)
                if hasattr(self, 'sp_cardapio'): self.sp_cardapio.config(from_=0, to=99, state='normal')


            if "TEF" not in fixed_quantities and "TEF" not in module_limits:
                self.spin_tef_var.set(0)
                if hasattr(self, 'sp_tef'): self.sp_tef.config(from_=0, to=99, state='normal')

            ##if "App Gestão CPlug" not in fixed_quantities and "App Gestão CPlug" not in module_limits:
            ##    self.spin_app_cplug_var.set(0)
            ##    if hasattr(self, 'sp_app_cplug'): self.sp_app_cplug.config(from_=0, to=999, state='normal')

            if "Delivery Direto Básico" not in fixed_quantities and "Delivery Direto Básico" not in module_limits:
                self.spin_delivery_direto_basico_var.set(0)
                if hasattr(self, 'sp_delivery_direto_basico'): self.sp_delivery_direto_basico.config(from_=0, to=999, state='normal')


            # --- Módulos Obrigatórios (Checkboxes) ---
            # Resetar todos os checkboxes de módulos opcionais e reabilitá-los
            for module_name, var_tk_int in self.modules.items():
                var_tk_int.set(0)
                if module_name in self.check_buttons:
                    self.check_buttons[module_name].config(state='normal')

            # Marcar e desabilitar os módulos obrigatórios do plano atual
            # CERTIFIQUE-SE QUE ESTA LINHA ESTÁ PRESENTE E CORRETAMENTE POSICIONADA:
            mandatory_modules_for_plan = info.get("mandatory", []) # <--- DEFINIÇÃO DA VARIÁVEL
     
            # Marcar e desabilitar os módulos obrigatórios do plano atual
            for mandatory_module in info.get("mandatory", []):
                if mandatory_module in self.modules:
                    self.modules[mandatory_module].set(1)
                    if mandatory_module in self.check_buttons:
                        self.check_buttons[mandatory_module].config(state='disabled')


            # --- Habilitar/Desabilitar Módulos Opcionais (Checkboxes) com base no Plano ---
            # Obter a lista de opcionais permitidos para este plano.
            # Se 'allowed_optionals' não estiver definido no PLAN_INFO para o plano,
            # assume-se que nenhum checkbox opcional é permitido (além dos obrigatórios).
            allowed_optionals_for_plan = info.get("allowed_optionals", []) # Default para lista vazia

            for module_name, var_tk_int in self.modules.items():
                if module_name in self.check_buttons:
                    # Ignorar módulos que já são obrigatórios para este plano (já foram tratados)
                    if module_name not in mandatory_modules_for_plan:
                        if module_name in allowed_optionals_for_plan:
                            # Mantém o estado 'normal' (já definido acima)
                            # Não desmarcar se já estiver marcado por alguma razão (embora deva estar 0 neste ponto)
                            pass
                        else:
                            var_tk_int.set(0) # Garante que está desmarcado se não for permitido
                            self.check_buttons[module_name].config(state='disabled')
          
                # else:
                #     print(f"Aviso: Módulo obrigatório '{mandatory_module}' não encontrado na lista de checkboxes self.modules.")

            # --- Lógica Específica de Notas Fiscais (do seu código original) ---
            # Tratar "3000 Notas Fiscais" como obrigatório para "Gestão (original)"
            # e "Notas Fiscais Ilimitadas" para outros, ou como parte de 'mandatory'.
            # A nova estrutura de PLAN_INFO já inclui isso em 'mandatory'.
            # A lógica abaixo pode precisar ser ajustada ou removida se 'mandatory' for suficiente.

            # Exemplo: Se "3000 Notas Fiscais" é mandatório para um plano específico
            if plano_selecionado == "Gestão (original)": # Use o nome exato do seu PLAN_INFO
                if "3000 Notas Fiscais" in self.modules:
                    self.modules["3000 Notas Fiscais"].set(1)
                    if "3000 Notas Fiscais" in self.check_buttons:
                        self.check_buttons["3000 Notas Fiscais"].config(state='disabled')
            else: # Para outros planos, desmarcar e habilitar "3000 Notas Fiscais" se não for mandatório
                if "3000 Notas Fiscais" in self.modules and "3000 Notas Fiscais" not in info.get("mandatory", []):
                    self.modules["3000 Notas Fiscais"].set(0)
                    if "3000 Notas Fiscais" in self.check_buttons:
                        self.check_buttons["3000 Notas Fiscais"].config(state='normal')


            # Resetar a seleção de notas por Radiobutton se o plano não for o PDV (ou conforme sua regra)
            # A descrição do Plano PDV tem "30 Notas Fiscais" como fixo. Os outros (60, 120, 250) são opcionais.
            # Para os planos Gestão e Performance, é "Notas Fiscais Ilimitadas".
            if plano_selecionado not in ["Plano PDV", "Plano PDV (original)"]: # Ajuste os nomes conforme seu PLAN_INFO
                self.var_notas.set("NONE") # "NONE" ou um valor que indique nenhuma seleção de radiobutton
                # Desabilitar radiobuttons se o plano já inclui "Notas Fiscais Ilimitadas"
                if "Notas Fiscais Ilimitadas" in info.get("mandatory", []):
                    # Aqui você precisaria iterar sobre os widgets Radiobutton e desabilitá-los
                    # Ex: for rb_widget in self.frame_left.winfo_children()[1].winfo_children()[0].winfo_children():
                    #         if isinstance(rb_widget, ttk.Radiobutton): rb_widget.config(state='disabled')
                    pass # Implemente a desabilitação dos Radiobuttons de notas aqui
                else:
                    # Habilitar Radiobuttons
                    # Ex: for rb_widget in ...: rb_widget.config(state='normal')
                    pass # Implemente a habilitação
            else: # Para Plano PDV
                self.var_notas.set("NONE") # Começa sem NF opcional selecionada
                # Habilitar Radiobuttons
                # Ex: for rb_widget in ...: rb_widget.config(state='normal')
                pass


            # --- Resetar Overrides de Cálculo ---
            self.user_override_anual_active.set(False)
            self.user_override_discount_active.set(False)
            self.valor_anual_editavel.set("0.00") # Será recalculado por atualizar_valores
            self.desconto_personalizado.set("0")  # Será recalculado

            # --- Atualizar Todos os Valores e a UI ---
            self.atualizar_valores()


    def atualizar_valores(self, *args):
            try:
                # Garante que self.current_plan foi definido antes de chamar atualizar_valores
                if not hasattr(self, 'current_plan') or not self.current_plan:
                    # Pode ser o caso durante a inicialização inicial antes de configurar_plano ser chamado.
                    # Defina um comportamento padrão ou retorne para evitar erros.
                    # print("Aviso: current_plan não definido em atualizar_valores. Usando plano padrão ou retornando.")
                    # self.lbl_plano_mensal.config(text="Plano (Mensal): R$ 0.00") # Exemplo de reset
                    # self.lbl_plano_anual.config(text="Plano (Anual): R$ 0.00")
                    # self.lbl_treinamento.config(text="Custo Treinamento (Mensal): R$ 0.00")
                    # self.lbl_desconto.config(text="Desconto: 0%")
                    return

                info = PLAN_INFO[self.current_plan]
            except KeyError:
                showerror("Erro de Cálculo", f"Plano '{self.current_plan}' não encontrado para cálculo.")
                return

            base_mensal_do_plano = info.get("base_mensal", 0.0)
            # 'parte_descontavel' começa com o valor base do plano, se este for descontável.
            # Se o base_mensal_do_plano em si não recebe desconto, ele iria para parte_sem_desc.
            # Assumindo que o valor base do plano É elegível para o desconto padrão anual.
            parte_descontavel = base_mensal_do_plano
            parte_sem_desc = 0.0

            mandatory_modules = info.get("mandatory", [])
            fixed_qty_modules = info.get("fixed_quantities", {}).keys() # Nomes dos módulos com quantidade fixa

            for module_name, var_tk_int in self.modules.items():
                if var_tk_int.get() == 1:
                    # Adiciona custo apenas se NÃO for obrigatório e NÃO for um módulo com quantidade fixa
                    # (cujo custo já deve estar no base_mensal_do_plano ou será tratado separadamente se não estiver)
                    if module_name not in mandatory_modules and module_name not in fixed_qty_modules:
                        
                        # Nova lógica para definir o preço do Backup Realtime
                        if module_name == "Backup Realtime":
                            if self.current_plan == "Performance":
                                module_price = 99.00  # Preço especial para o plano Performance
                            else:
                                # Busca o preço padrão em precos_mensais, ou usa 199.00 como fallback
                                module_price = precos_mensais.get("Backup Realtime", 199.00) 
                        else:
                            # Para todos os outros módulos, busca o preço em precos_mensais
                            module_price = precos_mensais.get(module_name, 0.0)

                        if module_name not in SEM_DESCONTO:
                            parte_descontavel += module_price
                        else:
                            parte_sem_desc += module_price

            # --- Notas Fiscais (Radiobuttons) ---
            # A lógica deve considerar se o plano já inclui um pacote de notas ou notas ilimitadas.
            # Se "Notas Fiscais Ilimitadas" ou "3000 Notas Fiscais" for mandatório,
            # os Radiobuttons de NF (60, 120, 250) não deveriam adicionar custo ou estar desabilitados.
            # Se o Plano PDV inclui "30 Notas Fiscais", o custo disso já está no base_mensal_do_plano.
            # As opções de Radiobutton seriam para upgrade.

            if "Notas Fiscais Ilimitadas" not in mandatory_modules and \
            "3000 Notas Fiscais" not in mandatory_modules and \
            "30 Notas Fiscais" not in mandatory_modules: # Verifica se NENHUM pacote de NF é mandatório
                # Se nenhum pacote de NF é mandatório, então o Radiobutton adiciona o custo total.
                selected_nf_option = self.var_notas.get() # Ex: "60", "120", "250", "NONE"
                if selected_nf_option != "NONE":
                    nf_module_name = f"{selected_nf_option} Notas Fiscais"
                    nf_price = precos_mensais.get(nf_module_name, 0.0)
                    if nf_module_name not in SEM_DESCONTO: # Assumindo que pacotes de NF podem ter desconto
                        parte_descontavel += nf_price
                    else:
                        parte_sem_desc += nf_price
            elif "30 Notas Fiscais" in mandatory_modules: # Ex: Plano PDV
                # O custo das 30 notas já está no base_mensal_do_plano.
                # Se o usuário selecionar um pacote MAIOR via Radiobutton, cobrar a DIFERENÇA ou o valor total do pacote maior.
                # Para simplificar, vamos assumir que o Radiobutton representa um UPGRADE e o preço é o do pacote selecionado.
                # Se desejar cobrar a diferença, a lógica precisaria do preço do pacote base.
                selected_nf_option = self.var_notas.get()
                if selected_nf_option != "NONE": # Usuário selecionou 60, 120, ou 250
                    nf_module_name = f"{selected_nf_option} Notas Fiscais"
                    nf_price = precos_mensais.get(nf_module_name, 0.0)
                    # Remova o custo do pacote base se o preço do radiobutton for o total e não a diferença
                    # parte_descontavel -= precos_mensais.get("30 Notas Fiscais", 0.0) # Se "30 Notas Fiscais" tiver um preço em precos_mensais
                    if nf_module_name not in SEM_DESCONTO:
                        parte_descontavel += nf_price
                    else:
                        parte_sem_desc += nf_price


            # --- PDVs Adicionais ---
            current_pdvs = self.spin_pdv_var.get()
            base_pdvs_no_plano = info.get("base_pdv", 0)
            max_additional_pdvs_allowed = info.get("max_additional_pdv", 0)
            
            pdvs_adicionais_reais = max(0, current_pdvs - base_pdvs_no_plano)
            pdvs_adicionais_cobradas = min(pdvs_adicionais_reais, max_additional_pdvs_allowed) # Não cobra além do permitido
            
            custo_pdv_adicional_unitario = info.get("pdv_extra_cost", 0.0)


            # Lógica específica para Bling (do seu código original, se ainda aplicável)
            if self.current_plan == "Bling (original)": # Ajuste o nome se necessário
                custo_pdv_adicional_unitario = info.get("pdv_extra_cost_bling", 40.00) # Exemplo
            
            # PDVs adicionais geralmente são descontáveis, mas verifique a regra de negócio.
            parte_descontavel += pdvs_adicionais_cobradas * custo_pdv_adicional_unitario

            # --- Usuários Adicionais ---
            current_users = self.spin_users_var.get()
            base_users_no_plano = info.get("base_users", 0)
            max_additional_users_allowed = info.get("max_additional_users", 0)

            users_adicionais_reais = max(0, current_users - base_users_no_plano)
            users_adicionais_cobradas = min(users_adicionais_reais, max_additional_users_allowed)

            custo_user_adicional_unitario = info.get("user_extra_cost", 20.00) # 20.00 era um default no seu código
            parte_descontavel += users_adicionais_cobradas * custo_user_adicional_unitario


            # --- Módulos Opcionais com Quantidade (Spinboxes) ---
            # TEF (geralmente sem desconto)
            parte_sem_desc += self.spin_tef_var.get() * precos_mensais.get("TEF", 99.00)

            # Smart TEF
            qty_smart_tef_selecionada = self.spin_smart_tef_var.get()
            if "Smart TEF" in fixed_qty_modules:
                # Custo dos Smart TEFs fixos já deve estar no base_mensal_do_plano.
                # Se não estiver, e precisar ser somado aqui:
                # custo_fixo_smart_tef = info["fixed_quantities"]["Smart TEF"] * precos_mensais.get("Smart TEF", 49.90)
                # if "Smart TEF" not in SEM_DESCONTO: parte_descontavel += custo_fixo_smart_tef
                # else: parte_sem_desc += custo_fixo_smart_tef
                pass # Assumindo que já está no base_mensal_do_plano
            elif "Smart TEF" in info.get("module_limits", {}): # Opcional com limite
                # O spinbox já está limitado, então qty_smart_tef_selecionada é válida.
                parte_sem_desc += qty_smart_tef_selecionada * precos_mensais.get("Smart TEF", 49.90)
            else: # Opcional sem limite específico para este plano (e não é fixo)
                if self.current_plan not in ["Plano Performance", "Plano Performance (original)"]: # Evitar dupla contagem
                    parte_sem_desc += qty_smart_tef_selecionada * precos_mensais.get("Smart TEF", 49.90)

            # Autoatendimento (lógica adaptada do seu original)
            auto_qty = self.spin_auto_var.get()
            if self.current_plan == "Autoatendimento": # Se for o PLANO Autoatendimento
                if auto_qty >= 1:
                    # O base_mensal_do_plano (R$0.0 no seu código) + base_anual (R$419.90) é confuso.
                    # Vamos assumir que o Plano Autoatendimento tem um custo base para o 1º terminal.
                    # Se o "base_anual" for o custo do primeiro, divida por 12 para mensal, ou use um "custo_base_terminal_aa".
                    # A sua lógica original somava base_anual aos custos se auto_qty >= 1. Isso precisa ser esclarecido.
                    # Para este exemplo, vamos seguir a lógica de preço de módulo opcional.
                    # Se o "Plano Autoatendimento" tiver um preço fixo para o primeiro terminal, isso deve estar no base_mensal_do_plano.
                    # E então os adicionais teriam um custo diferente.
                    # A descrição textual não tem um "Plano Autoatendimento" mas sim "Terminais Autoatendimento (R$ 199,00)" como opcional.
                    # Vou usar o preço de módulo opcional "Autoatendimento" de precos_mensais.
                    parte_sem_desc += auto_qty * precos_mensais.get("Autoatendimento", 299.90) # R$299.90 era o preço do módulo opcional
            else: # Autoatendimento como módulo opcional em outros planos
                parte_sem_desc += auto_qty * precos_mensais.get("Autoatendimento", 299.90)


            # App Gestão CPlug (descontável)
            # parte_descontavel += self.spin_app_cplug_var.get() * precos_mensais.get("App Gestão CPlug", 19.90)

            # Delivery Direto Básico (descontável)
            parte_descontavel += self.spin_delivery_direto_basico_var.get() * precos_mensais.get("Delivery Direto Básico", 247.00)

            # Cardápio Digital (descontável, com lógica de preço por quantidade)
            card_qt = self.spin_cardapio_var.get()
            if card_qt == 1:
                parte_descontavel += precos_mensais.get("Cardápio Digital_unitario", 29.90) # Preço para 1 unidade
            elif card_qt > 1:
                # Se o preço unitário para múltiplos for diferente.
                # Assumindo que "Cardápio Digital" em precos_mensais é o preço para o primeiro,
                # e "Cardápio Digital_multiplo" é para os seguintes, ou um preço fixo por unidade se > 1.
                # Sua lógica original: card_qt * 24.90.
                # Vamos usar um preço de precos_mensais para múltiplos.
                preco_multiplo = precos_mensais.get("Cardápio Digital_multiplo", 24.90)
                # Se o primeiro tem preço diferente e os demais outro:
                # parte_descontavel += precos_mensais.get("Cardápio Digital_unitario", 29.90) + (card_qt - 1) * preco_multiplo
                # Se todos têm o mesmo preço quando > 1:
                parte_descontavel += card_qt * preco_multiplo


            # --- Cálculo Final dos Valores ---
            valor_mensal_automatico = parte_descontavel + parte_sem_desc

            # Cálculo Anual e Desconto
            final_anual = 0.0
            if self.user_override_anual_active.get():
                try:
                    final_anual = float(self.valor_anual_editavel.get())
                except ValueError:
                    final_anual = valor_mensal_automatico # Fallback se a entrada for inválida
                    self.valor_anual_editavel.set(f"{final_anual:.2f}")
            elif self.user_override_discount_active.get():
                try:
                    desc_custom_percent = float(self.desconto_personalizado.get())
                    if not (0 <= desc_custom_percent <= 100): # Valida percentual
                        desc_custom_percent = 0.0
                        self.desconto_personalizado.set("0")
                except ValueError:
                    desc_custom_percent = 0.0
                    self.desconto_personalizado.set("0")
                
                desc_decimal = desc_custom_percent / 100.0
                final_anual = (parte_descontavel * (1 - desc_decimal)) + parte_sem_desc # Desconto só na parte descontável
                self.valor_anual_editavel.set(f"{final_anual:.2f}")
            else: # Desconto padrão de 10% (ou outro valor que você definir)
                desc_padrao_decimal = 0.10 # 10%
                final_anual = (parte_descontavel * (1 - desc_padrao_decimal)) + parte_sem_desc
                self.valor_anual_editavel.set(f"{final_anual:.2f}") # Atualiza o campo editável com o cálculo automático


            # Custo Treinamento (lógica do seu código original)
            training_cost = 0.0
            if self.current_plan == "Autoatendimento" or self.current_plan == "Em Branco":
                training_cost = 0.0
            else:
                if valor_mensal_automatico < 499.00: # Limiar para custo de treinamento
                    training_cost = 499.00 - valor_mensal_automatico
                else:
                    training_cost = 0.0

            # --- Atualização das Labels na UI ---
            if self.spin_auto_var.get() > 0 and self.current_plan != "Autoatendimento": # Se for módulo opcional
                self.lbl_plano_mensal.config(text=f"Plano (Mensal): R$ {valor_mensal_automatico:.2f}")
            elif self.current_plan == "Autoatendimento": # Regra específica do plano autoatendimento
                # A sua lógica original mostrava "Não disponível" para mensal no plano Autoatendimento.
                # Se o "base_anual" for o único preço, o mensal seria base_anual / 12 ou não aplicável.
                # Se o plano Autoatendimento tiver um base_mensal em PLAN_INFO, use-o.
                # Vamos assumir que o valor_mensal_automatico já reflete isso.
                # Se o plano Autoatendimento for APENAS anual, então mensal seria "Não disponível".
                self.lbl_plano_mensal.config(text="Plano (Mensal): Não disponível") # Ou o cálculo mensal se existir
            else:
                self.lbl_plano_mensal.config(text=f"Plano (Mensal): R$ {valor_mensal_automatico:.2f}")

            self.lbl_plano_anual.config(text=f"Plano (Anual): R$ {final_anual:.2f}")
            self.lbl_treinamento.config(text=f"Custo Treinamento (Mensal): R$ {training_cost:.2f}")



# ... (valor_mensal_automatico e final_anual são determinados com base em toda a lógica, incluindo overrides) ...

            # NOVA LÓGICA PARA DESCONTO EXIBIDO E ARMAZENADO
            effective_discount_percent = 0.0
            if valor_mensal_automatico > 0:
                # Calcula a economia monetária total
                total_saving = valor_mensal_automatico - final_anual
                
                # Calcula o percentual de desconto efetivo com base na economia total e no valor mensal total
                # Permite que total_saving seja negativo se final_anual for definido manualmente como um valor alto.
                effective_discount_percent = (total_saving / valor_mensal_automatico) * 100
            # Se valor_mensal_automatico for 0 ou negativo, effective_discount_percent permanece 0.0

            # Atualiza a etiqueta da UI com o desconto efetivo arredondado
            self.lbl_desconto.config(text=f"Desconto Aplicado: {round(effective_discount_percent)}%")

            # Armazenar valores computados para uso na geração da proposta
            self.computed_mensal = valor_mensal_automatico
            self.computed_anual = final_anual
            # Armazena o percentual de desconto efetivo real
            self.computed_desconto_percent = round(effective_discount_percent) 

            # --- Calcular Custo de Setup Adicional para a UI e Proposta ---
            custo_total_setup_val = 0.0
            for nome_modulo, var_modulo_selecionado in self.modules.items():
                if var_modulo_selecionado.get() == 1:
                    if nome_modulo in CUSTOS_SETUP_ADICIONAL:
                        custo_total_setup_val += CUSTOS_SETUP_ADICIONAL[nome_modulo]

            self.computed_custo_setup_adicional_valor = custo_total_setup_val # Armazena o valor numérico

            texto_label_setup = "Custo Adicional Setup: R$ 0,00"
            string_placeholder_setup = ""

            if custo_total_setup_val > 0:
                valor_formatado_ui = f"{custo_total_setup_val:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
                texto_label_setup = f"Custo Adicional Setup: R$ {valor_formatado_ui}"
                string_placeholder_setup = f"+ Setup R$ {valor_formatado_ui}" # A mesma formatação para o placeholder

            if hasattr(self, 'lbl_custo_setup_adicional'): # Garante que o label existe
                self.lbl_custo_setup_adicional.config(text=texto_label_setup)

            self.computed_custo_setup_adicional_str = string_placeholder_setup # Armazena a string para o placeholder da proposta


    def montar_lista_modulos(self):
            # Garante que self.current_plan e PLAN_INFO estão disponíveis
            if not hasattr(self, 'current_plan') or not self.current_plan or not PLAN_INFO:
                return ["Erro: Informações do plano não disponíveis para montar a lista de módulos."]
            try:
                info = PLAN_INFO[self.current_plan]
            except KeyError:
                return [f"Erro: Plano '{self.current_plan}' não encontrado em PLAN_INFO."]

            linhas_quantidades = []     # Para itens como "Nx Módulo"
            linhas_checkboxes_etc = []  # Para itens como " Módulo"
            modulos_ja_listados = set() # Para evitar duplicar módulos na lista

            # --- SEÇÃO A: Módulos com quantidade explícita ---

            # A.1: PDVs (Consolidado) e Usuário(s) Cortesia
            total_pdvs_selecionados = self.spin_pdv_var.get() # Total de PDVs que o usuário configurou
            base_pdvs_no_plano = info.get("base_pdv", 0)
            
            # Calcula quantos PDVs são ADICIONAIS em relação à base do plano
            pdvs_adicionais_comprados = max(0, total_pdvs_selecionados - base_pdvs_no_plano)

            # Adiciona a linha consolidada de PDVs
            if total_pdvs_selecionados > 0:
                texto_pdv_consolidado = f"{total_pdvs_selecionados}x PDV - Frente de Caixa "
                # Adiciona à lista se o nome base do módulo ainda não foi registrado
                # (evita duplicidade caso "PDV - Frente de Caixa" viesse de fixed_quantities, por exemplo, o que é improvável aqui)
                if "PDV - Frente de Caixa" not in modulos_ja_listados:
                    linhas_quantidades.append(texto_pdv_consolidado)
                    modulos_ja_listados.add("PDV - Frente de Caixa") # Registra o nome base do módulo

            # Adiciona "Usuário(s) Cortesia" com base nos PDVs ADICIONAIS
            if pdvs_adicionais_comprados > 0:
                if pdvs_adicionais_comprados == 1:
                    linhas_quantidades.append(f"1x Usuário Cortesia ")
                else: # pdvs_adicionais_comprados > 1
                    linhas_quantidades.append(f"{pdvs_adicionais_comprados}x Usuários Cortesia ")
                # Não é necessário adicionar "Usuário Cortesia" ou "Usuários Cortesia" a modulos_ja_listados,
                # pois são itens de cortesia distintos e não módulos que podem ser selecionados em outro lugar.

            # A.2: Usuários (Base e Adicionais)
            current_users = self.spin_users_var.get() # Total de usuários que o usuário configurou
            base_users_no_plano = info.get("base_users", 0)
            users_adicionais_comprados = max(0, current_users - base_users_no_plano)

            # Adiciona a linha de usuários base (ou total, se não houver distinção na exibição)
            if current_users > 0: # Mostra o total de usuários configurados
                # Se a intenção é mostrar usuários base + adicionais separadamente, a lógica precisaria mudar.
                # Para mostrar o total de usuários configurados em uma linha:
                texto_user_total = f"{current_users}x Usuários "
                if "Usuários" not in modulos_ja_listados: # Verifica se "Usuários" (como conceito geral) já foi listado
                    linhas_quantidades.append(texto_user_total)
                    modulos_ja_listados.add("Usuários") # Adiciona o nome base
                # Se você quisesse listar "Usuário Adicional" separadamente, a lógica original seria:
                # if base_users_no_plano > 0:
                #     texto_user = f"{base_users_no_plano}x Usuários " # Ou o nome do módulo de usuários base
                #     if "Usuários" not in modulos_ja_listados: # Ou nome do módulo de usuários base
                #         linhas_quantidades.append(texto_user)
                #         modulos_ja_listados.add("Usuários") # Ou nome do módulo de usuários base
                # if users_adicionais_comprados > 0:
                #     texto_user_extra = f"{users_adicionais_comprados}x Usuário Adicional "
                #     if "Usuário Adicional" not in modulos_ja_listados:
                #         linhas_quantidades.append(texto_user_extra)
                #         modulos_ja_listados.add("Usuário Adicional")
            
            # A.3: Módulos com Quantidade Fixa (Definidos em 'fixed_quantities')
            fixed_quantities = info.get("fixed_quantities", {})
            for mod_fixo, qty_fixa in fixed_quantities.items():
                if mod_fixo not in modulos_ja_listados:
                    linhas_quantidades.append(f"{qty_fixa}x {mod_fixo} ")
                    modulos_ja_listados.add(mod_fixo)

            # A.4: Módulos Opcionais com Quantidade (Controlados por Spinbox)
            # TEF
            qty_tef = self.spin_tef_var.get()
            if qty_tef > 0 and "TEF" not in fixed_quantities and "TEF" not in modulos_ja_listados:
                linhas_quantidades.append(f"{qty_tef}x TEF ")
                modulos_ja_listados.add("TEF")

            # Smart TEF
            qty_smart_tef = self.spin_smart_tef_var.get()
            if "Smart TEF" not in fixed_quantities and qty_smart_tef > 0 and "Smart TEF" not in modulos_ja_listados:
                linhas_quantidades.append(f"{qty_smart_tef}x Smart TEF ")
                modulos_ja_listados.add("Smart TEF")

            # Autoatendimento
            qty_auto = self.spin_auto_var.get()
            if "Autoatendimento" not in fixed_quantities and qty_auto > 0 and "Autoatendimento" not in modulos_ja_listados:
                linhas_quantidades.append(f"{qty_auto}x Terminal de Autoatendimento ")
                modulos_ja_listados.add("Autoatendimento") 

                linhas_quantidades.append(f"{qty_auto}x TEF Cortesia ")
            
            # --- SEÇÃO B: Módulos estilo checkbox e Notas Fiscais opcionais ---

            # B.1: Módulos Obrigatórios (Definidos em 'mandatory')
            mandatory_modules = info.get("mandatory", [])
            for obrig_module in mandatory_modules:
                if obrig_module not in modulos_ja_listados:
                    linhas_checkboxes_etc.append(f" {obrig_module} ")
                    modulos_ja_listados.add(obrig_module)

            # B.2: Notas Fiscais (Radiobuttons - Opcionais/Upgrade)
            selected_nf_option = self.var_notas.get()
            if selected_nf_option != "NONE":
                nf_module_name_selected = f"{selected_nf_option} Notas Fiscais"
                if "Notas Fiscais Ilimitadas" not in mandatory_modules and \
                nf_module_name_selected not in modulos_ja_listados:
                    linhas_checkboxes_etc.append(f" {nf_module_name_selected} ")
                    modulos_ja_listados.add(nf_module_name_selected)

            # B.3: Módulos Opcionais (Checkboxes Selecionados de self.modules)
            for module_name, var_tk_int in self.modules.items():
                if var_tk_int.get() == 1: 
                    if module_name not in modulos_ja_listados:
                        linhas_checkboxes_etc.append(f" {module_name} ")
                        modulos_ja_listados.add(module_name)

            # --- Combina as listas na ordem desejada ---
            linhas_proposta = linhas_quantidades + linhas_checkboxes_etc
            return linhas_proposta



    def gerar_dados_proposta(self, nome_closer, cel_closer, email_closer):
            nome_plano = self.nome_plano_var.get().strip() or "Plano"

            valor_mensal = self.computed_mensal
            valor_anual = self.computed_anual
            desconto_percent = self.computed_desconto_percent

            auto_qty = self.spin_auto_var.get()
            if auto_qty > 0:
                plano_mensal_str = "Não Disponível"
                training_cost = 0.0
            else:
                training_cost = 0.0
                if valor_mensal < 499.00:
                    training_cost = 499.00 - valor_mensal
                if training_cost > 0:
                    part_mensal = f"{valor_mensal:.2f}".replace(".", ",")
                    part_training = f"{training_cost:.2f}".replace(".", ",")
                    plano_mensal_str = f"R$ {part_mensal} + R$ {part_training}"
                else:
                    plano_mensal_str = f"R$ {valor_mensal:.2f}".replace(".", ",")

            plano_anual_str = f"R$ {valor_anual:.2f}".replace(".", ",")

            if "Suporte Técnico - Estendido" in self.modules and self.modules["Suporte Técnico - Estendido"].get() == 1:
                tipo_suporte = "Estendido"
                horario_suporte = "09:00 às 22:00 de Segunda a Sexta-feira & Sábado e Domingo das 11:00 às 21:00"
            else:
                tipo_suporte = "Regular"
                horario_suporte = "09:00 às 17:00 de Segunda a Sexta-feira"

        
            lista_mods_formatada = self.montar_lista_modulos() # Retorna strings já formatadas para o conteúdo


            # --- AJUSTE DA LISTA DE MÓDULOS ---
            
            lista_mods_conteudo = self.montar_lista_modulos() # Retorna strings de conteúdo dos módulos

            LIMITE_LINHAS_EXTENSAO = 18 
            
            # Adiciona o bullet point a cada item da lista de conteúdo
            linhas_com_bullet = [f"• {m}" for m in lista_mods_conteudo]

            montagem_plano_principal_str = ""
            extensao_plano_str = "" # Novo nome do placeholder para a extensão

            if len(linhas_com_bullet) > LIMITE_LINHAS_EXTENSAO:
                # Se a lista total tem MAIS de 18 linhas:
                # As primeiras 18 linhas vão para "extecao_do_plano"
                extensao_plano_str = "\n".join(linhas_com_bullet[:LIMITE_LINHAS_EXTENSAO])
                # O restante (da 19ª linha em diante) vai para "montagem_do_plano"
                montagem_plano_principal_str = "\n".join(linhas_com_bullet[LIMITE_LINHAS_EXTENSAO:])
            else:
                # Se a lista total tem 18 linhas OU MENOS:
                # Todas as linhas vão para "montagem_do_plano"
                montagem_plano_principal_str = "\n".join(linhas_com_bullet)
                # "extecao_do_plano" fica vazia
                extensao_plano_str = ""
            # --- FIM DO AJUSTE DA LISTA DE MÓDULOS ---


            lista_mods = self.montar_lista_modulos()
            montagem = "\n".join(f"•    {m}" for m in lista_mods)

            if self.current_plan == "Autoatendimento":
                economia_str = ""
            else:
                custo_anual_mensalizado = valor_mensal * 12 + training_cost
                custo_anual_plano = valor_anual * 12
                economia_val = custo_anual_mensalizado - custo_anual_plano
                if economia_val > 0:
                    econ = f"{economia_val:.2f}".replace(".", ",")
                    economia_str = f"Economia de R$ {econ} ao ano"
                else:
                    economia_str = ""


            # --- INÍCIO DO CÁLCULO DO SETUP ADICIONAL ---
            custo_total_setup = 0.0
            for nome_modulo, var_modulo_selecionado in self.modules.items():
                if var_modulo_selecionado.get() == 1: # Verifica se o módulo (checkbox) está selecionado
                    if nome_modulo in CUSTOS_SETUP_ADICIONAL:
                        custo_total_setup += CUSTOS_SETUP_ADICIONAL[nome_modulo]
            
            string_setup_adicional = ""
            if custo_total_setup > 0:
                # Formata para o padrão R$ 1.000,00
                valor_formatado_setup = f"{custo_total_setup:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
                string_setup_adicional = f"+ Setup R$ {valor_formatado_setup}"

            dados = {
                "montagem_do_plano": montagem_plano_principal_str, 
                "extecao_do_plano": extensao_plano_str, 
                "plano_mensal": plano_mensal_str,
                "plano_anual": plano_anual_str,
                "desconto_total": f"{desconto_percent}%",
                "nome_do_plano": nome_plano,
                "tipo_de_suporte": tipo_suporte,
                "horario_de_suporte": horario_suporte,
                "validade_proposta": self.validade_proposta_var.get(),
                "nome_closer": nome_closer,
                "celular_closer": cel_closer,
                "email_closer": email_closer,
                "nome_cliente": self.nome_cliente_var.get(),
                "economia_anual": economia_str,
                "setup_adicional": string_setup_adicional
            }

            return dados


# ---------------------------------------------------------
# Funções que geram .pptx (Proposta e Material)
# ---------------------------------------------------------
def gerar_proposta(lista_abas, nome_closer, celular_closer, email_closer):
    ppt_file = "Proposta Comercial ConnectPlug.pptx"
    if not os.path.exists(ppt_file):
        showerror("Erro", f"Arquivo '{ppt_file}' não encontrado!")
        return None

    try:
        prs = Presentation(ppt_file)
    except Exception as e:
        showerror("Erro", f"Falha ao abrir '{ppt_file}': {e}")
        return None

    if not lista_abas:
        showerror("Erro", "Não há abas para gerar Proposta.")
        return None

    # 1) Descobrir quais slides manter (opcional)
    abas_indices = sorted([aba.aba_index for aba in lista_abas])
    used_plans = {aba.current_plan for aba in lista_abas}
    
    keep_slides = set()
    slide_map_aba = {}
    
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texts.append(run.text)
        full_txt = " ".join(texts)

        # Exemplo: se tiver "slide_bling", só mantém se "Bling" estiver em used_plans
        if "slide_bling" in full_txt:
            if "Bling" not in used_plans:
                continue
        
        found_aba = None
        if "aba_plano_" not in full_txt:
            # Slide genérico => mantém
            keep_slides.add(i)
            slide_map_aba[i] = None
        else:
            # Slide para aba específica
            for x in abas_indices:
                marker = f"aba_plano_{x}"
                if marker in full_txt:
                    found_aba = x
                    break
            if found_aba is not None:
                keep_slides.add(i)
                slide_map_aba[i] = found_aba
    
    # 2) Remover slides não mantidos
    for idx in reversed(range(len(prs.slides))):
        if idx not in keep_slides:
            rid = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rid)
            del prs.slides._sldIdLst[idx]
    
    # 3) Re-mapear índices
    sorted_kept = sorted(keep_slides)
    new_order_map = {}
    for new_idx, slide in enumerate(prs.slides):
        old_idx = sorted_kept[new_idx]
        new_order_map[new_idx] = slide_map_aba[old_idx]
    
    # 4) Substituir placeholders
    dados_de_aba = {}
    for aba in lista_abas:
        d = aba.gerar_dados_proposta(nome_closer, celular_closer, email_closer)
        dados_de_aba[aba.aba_index] = d
    
    # Fallback: se não tiver slides específicos, use dados da primeira aba
    fallback_aba = lista_abas[0]
    d_fallback = dados_de_aba[fallback_aba.aba_index]

    for new_idx, slide in enumerate(prs.slides):
        aba_num = new_order_map[new_idx]
        if aba_num is None:
            substituir_placeholders_no_slide(slide, d_fallback)
        else:
            d_aba = dados_de_aba[aba_num]
            substituir_placeholders_no_slide(slide, d_aba)
    
    # 5) Salvar
    nome_cliente_primeira = d_fallback.get("nome_cliente", "SemNome")
    hoje_str = date.today().strftime("%d-%m-%Y")
    nome_arquivo = f"Proposta ConnectPlug - {nome_cliente_primeira} - {hoje_str}.pptx"

    try:
        prs.save(nome_arquivo)
        showinfo("Sucesso", f"Proposta gerada: {nome_arquivo}")
        return nome_arquivo
    except Exception as e:
        showerror("Erro", f"Falha ao salvar: {e}")
        return None

def gerar_material(lista_abas, nome_closer, celular_closer, email_closer):
    mat_file = "Material Tecnico ConnectPlug.pptx"
    if not os.path.exists(mat_file):
        showerror("Erro", f"Arquivo '{mat_file}' não encontrado!")
        return None

    try:
        prs = Presentation(mat_file)
    except Exception as e:
        showerror("Erro", f"Falha ao abrir '{mat_file}': {e}")
        return None

    if not lista_abas:
        showerror("Erro", "Não há abas para gerar Material Técnico.")
        return None

    # ---------------------------------------------------
    # 1) Descobrir módulos ativos e planos usados
    # ---------------------------------------------------
    modulos_ativos = set()
    planos_usados = set()

    for aba in lista_abas:
        planos_usados.add(aba.current_plan)

        # Módulos de checkboxes
        for nome_mod, var_mod in aba.modules.items():
            if var_mod.get() == 1:
                modulos_ativos.add(nome_mod)

        # Incrementos
        if aba.spin_tef_var.get() > 0:
            modulos_ativos.add("TEF")
        if aba.spin_smart_tef_var.get() > 0:
            modulos_ativos.add("Smart TEF")
        if aba.spin_auto_var.get() > 0:
            modulos_ativos.add("Autoatendimento")
        if aba.spin_cardapio_var.get() > 0:
            modulos_ativos.add("Cardápio Digital")
        #if aba.spin_app_cplug_var.get() > 0:
        #    modulos_ativos.add("App Gestão CPlug")
        #if aba.spin_delivery_direto_basico_var.get() > 0:
        #    modulos_ativos.add("Delivery Direto Básico")
        if aba.spin_pdv_var.get() > 0:
            modulos_ativos.add("PDV")
        # etc., se houver outros increments.

    # ---------------------------------------------------
    # 2) Mapeamento de placeholders para módulos
    #    Adapte conforme seus slides
    # ---------------------------------------------------
    # Se um slide contiver o texto "check_tef" e você quiser mantê-lo só se
    # "TEF" estiver em modulos_ativos, defina assim:
    MAPEAMENTO_MODULOS = {
        "slide_sempre": None,
        "check_sistema_kds": "Relatório KDS",
        "check_Hub_de_Delivery": "Hub de Delivery",
        "check_integracao_api": "Integração API",
        "check_integracao_tap": "Integração TAP",
        "check_controle_de_mesas": "Controle de Mesas",
        "check_Delivery": "Delivery",
        "check_producao": "Produção",
        "check_Estoque_em_Grade": "Estoque em Grade",
        "check_Facilita_NFE": "Facilita NFE",
        "check_Importacao_de_xml": "Importação de XML",
        "check_conciliacao_bancaria": "Conciliação Bancária",
        "check_contratos_de_cartoes": "Contratos de cartões e outros",
        "check_ordem_de_servico": "Ordem de Serviço",
        "check_relatorio_dinamico": "Relatório Dinâmico",
        "check_programa_de_fidelidade": "Programa de Fidelidade",
        "check_business_intelligence": "Business Intelligence (BI)",
        "check_smartmenu": "Smart Menu",
        "check_backup_real_time": "Backup Realtime",
        "check_att_tempo_real": "Atualização em Tempo Real",
        "check_promocao": "Promoções",
        "check_marketing": "Marketing",
        "pdv_balcao": "PDV",
        "qtd_smarttef": "Smart TEF",
        "qtd_tef": "TEF",
        "qtd_autoatendimento": "Autoatendimento",
        "qtd_cardapio_digital": "Cardápio Digital",
        "qtd_app_gestao_cplug": "App Gestão CPlug",
        "qtd_delivery_direto_basico": "Delivery Direto Básico",
        "check_delivery_direto_vip": "Delivery Direto VIP",
        "check_delivery_direto_profissional": "Delivery Direto Profissional",
        "check_notas_fiscais": {
            "Notas Fiscais Ilimitadas",
            "60 Notas Fiscais",
            "120 Notas Fiscais",
            "250 Notas Fiscais",
            "3000 Notas Fiscais"
        },
    }

    # ---------------------------------------------------
    # 3) Decidir quais slides manter
    # ---------------------------------------------------
    keep_slides = set()

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        texts.append(run.text.strip())
        full_txt = " ".join(texts)

        # Flag para saber se iremos manter este slide
        slide_ok = False

        # Caso tenha "slide_bling" e Bling esteja em use, etc.
        if "slide_bling" in full_txt:
            if "Bling" in planos_usados:
                slide_ok = True

        # Se tiver "slide_sempre" => mantém sempre
        if "slide_sempre" in full_txt:
            slide_ok = True

        # Agora checamos placeholders de módulos do MAPEAMENTO_MODULOS
        for placeholder, nome_modulo in MAPEAMENTO_MODULOS.items():
            if placeholder in full_txt:
                # Se o placeholder está no slide, mantenha só se
                # esse nome_modulo estiver em modulos_ativos
                if nome_modulo in modulos_ativos:
                    slide_ok = True
                    # Não fazemos break aqui, pois podem haver vários placeholders
                    # no mesmo slide. Mas se quiser, pode fazer break se
                    # bastar um para manter.

        # Exemplo: Se não entrou em nenhum if e slide_ok ainda está False,
        # esse slide NÃO será mantido.
        if slide_ok:
            keep_slides.add(i)

    # 4) Remove slides não mantidos
    for idx in reversed(range(len(prs.slides))):
        if idx not in keep_slides:
            rid = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rid)
            del prs.slides._sldIdLst[idx]

    # ---------------------------------------------------
    # 5) Substituir placeholders (dados globais)
    # ---------------------------------------------------
    fallback_aba = lista_abas[0]
    d_fallback = fallback_aba.gerar_dados_proposta(nome_closer, celular_closer, email_closer)

    for slide in prs.slides:
        substituir_placeholders_no_slide(slide, d_fallback)

    # ---------------------------------------------------
    # 6) Salvar pptx final
    # ---------------------------------------------------
    nome_cliente_primeira = d_fallback.get("nome_cliente", "SemNome")
    hoje_str = date.today().strftime("%d-%m-%Y")
    nome_arquivo = f"Material Tecnico ConnectPlug - {nome_cliente_primeira} - {hoje_str}.pptx"

    try:
        prs.save(nome_arquivo)
        showinfo("Sucesso", f"Material Técnico gerado: {nome_arquivo}")
        return nome_arquivo
    except Exception as e:
        showerror("Erro", f"Falha ao salvar: {e}")
        return None



# ---------------------------------------------------------
# Google Drive / Auth
# ---------------------------------------------------------
SCOPES = ['https://www.googleapis.com/auth/drive']

def baixar_client_secret_remoto():
    """Baixa o client_secret.json do repositório no GitHub se ainda não estiver salvo localmente."""
    url = "https://github.com/DevRGS/Gerador/raw/refs/heads/main/config/client_secret_788265418970-ur6f189oqvsttseeg6g77fegt0su67dj.apps.googleusercontent.com.json"
    nome_local = "client_secret_temp.json"

    if not os.path.exists(nome_local):
        print("Baixando client_secret do GitHub...")
        r = requests.get(url)
        if r.status_code == 200:
            with open(nome_local, "w", encoding="utf-8") as f:
                f.write(r.text)
        else:
            raise Exception(f"Erro ao baixar o client_secret.json: {r.status_code}")
    
    return nome_local

def get_gdrive_service():
    """Autentica e retorna o serviço do Google Drive com base no client_secret remoto."""
    creds = None
    CLIENT_SECRET_FILE = baixar_client_secret_remoto()
    TOKEN_FILE = 'token.json'

    # Tenta carregar o token salvo anteriormente
    if os.path.exists(TOKEN_FILE):
        with open(TOKEN_FILE, 'rb') as token:
            creds = pickle.load(token)

    # Se não houver token ou ele for inválido/expirado, roda o fluxo de autenticação
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CLIENT_SECRET_FILE, SCOPES)
            creds = flow.run_local_server(port=0)

        # Salva o token local para reutilização
        with open(TOKEN_FILE, 'wb') as token:
            pickle.dump(creds, token)

    # Constrói o serviço Google Drive autenticado
    service = build('drive', 'v3', credentials=creds)
    return service

def upload_pptx_and_export_to_pdf(local_pptx_path):
    """
    Faz upload do .pptx convertendo em Google Slides, 
    e baixa PDF local trocando .pptx -> .pdf
    """
    if not os.path.exists(local_pptx_path):
        showerror("Erro", f"Arquivo {local_pptx_path} não foi encontrado.")
        return

    service = get_gdrive_service()
    pdf_output_name = local_pptx_path.replace(".pptx", ".pdf")

    # 1) Upload (convertendo)
    file_metadata = {
        'name': os.path.basename(local_pptx_path),
        'mimeType': 'application/vnd.google-apps.presentation'
    }
    media = MediaFileUpload(
        local_pptx_path,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        resumable=True
    )
    uploaded_file = service.files().create(
        body=file_metadata,
        media_body=media,
        fields='id'
    ).execute()
    file_id = uploaded_file.get('id')
    print(f"Arquivo '{local_pptx_path}' enviado como Google Slides. ID: {file_id}")

    # 2) Exportar para PDF
    request = service.files().export_media(fileId=file_id, mimeType='application/pdf')
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        status, done = downloader.next_chunk()
        if status:
            print(f"Progresso PDF: {int(status.progress() * 100)}%")

    with open(pdf_output_name, 'wb') as f:
        f.write(fh.getvalue())

    showinfo("Google Drive", f"PDF gerado localmente: '{pdf_output_name}'.")


# ---------------------------------------------------------
# MainApp
# ---------------------------------------------------------
class MainApp(ttkb.Window):
    def __init__(self):
        super().__init__(themename="litera")
        self.title("Proposta + Material Técnico + PDF unificado")
        self.geometry("1200x800")

        self.nome_closer_var = tk.StringVar()
        self.celular_closer_var = tk.StringVar()
        self.email_closer_var = tk.StringVar()

        # Variáveis compartilhadas para TODAS as abas
        self.nome_cliente_var_shared = tk.StringVar(value="Cliente Compartilhado")
        self.validade_proposta_var_shared = tk.StringVar(value="DD/MM/YYYY")

        carregar_config(self.nome_closer_var, self.celular_closer_var, self.email_closer_var)
        self.protocol("WM_DELETE_WINDOW", self.on_close)

        # Top bar
        top_bar = ttkb.Frame(self)
        top_bar.pack(side="top", fill="x", pady=5)

        ttkb.Label(top_bar, text="Vendedor:").pack(side="left", padx=5)
        ttkb.Entry(top_bar, textvariable=self.nome_closer_var, width=15).pack(side="left", padx=5)
        ttkb.Label(top_bar, text="Cel:").pack(side="left", padx=5)
        ttkb.Entry(top_bar, textvariable=self.celular_closer_var, width=15).pack(side="left", padx=5)
        ttkb.Label(top_bar, text="Email:").pack(side="left", padx=5)
        ttkb.Entry(top_bar, textvariable=self.email_closer_var, width=20).pack(side="left", padx=5)

        self.btn_add = ttkb.Button(top_bar, text="+ Nova Aba", command=self.add_aba)
        self.btn_add.pack(side="right", padx=5)

        self.notebook = ttkb.Notebook(self)
        self.notebook.pack(fill="both", expand=True)

        bot_frame = ttkb.Frame(self)
        bot_frame.pack(side="bottom", fill="x", pady=5)

        # Botões unificados (gera .pptx e PDF)
        ttkb.Button(bot_frame, text="Gerar Proposta + PDF", command=self.on_gerar_proposta).pack(side="left", padx=5)
        ttkb.Button(bot_frame, text="Gerar Material + PDF", command=self.on_gerar_mat_tecnico).pack(side="left", padx=5)
        ttkb.Button(bot_frame, text="Gerar TUDO + PDF", command=self.on_gerar_tudo).pack(side="left", padx=5)

        self.abas_criadas = {}
        self.ultimo_indice = 0

        # Cria ao menos 1 aba no começo
        self.add_aba()
        baixar_arquivo_if_needed(
        "Proposta Comercial ConnectPlug.pptx",
        "https://github.com/DevRGS/Gerador/raw/refs/heads/main/assets/Proposta%20Comercial%20ConnectPlug.pptx")

        baixar_arquivo_if_needed(
        "Material Tecnico ConnectPlug.pptx",
        "https://github.com/DevRGS/Gerador/raw/refs/heads/main/assets/Material%20Tecnico%20ConnectPlug.pptx" )

    def on_close(self):
        salvar_config(
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        self.destroy()

    def add_aba(self):
        if len(self.abas_criadas) >= MAX_ABAS:
            return
        self.ultimo_indice += 1
        idx = self.ultimo_indice
        frame_aba = PlanoFrame(
            self.notebook,
            idx,
            nome_cliente_var_shared=self.nome_cliente_var_shared,
            validade_proposta_var_shared=self.validade_proposta_var_shared,
            on_close_callback=self.fechar_aba
        )
        self.notebook.add(frame_aba, text=f"Aba {idx}")
        self.abas_criadas[idx] = frame_aba

    def fechar_aba(self, indice):
        if indice in self.abas_criadas:
            frame_aba = self.abas_criadas[indice]
            self.notebook.forget(frame_aba)
            del self.abas_criadas[indice]

    def get_abas_ativas(self):
        return [self.abas_criadas[k] for k in sorted(self.abas_criadas.keys())]

    def on_gerar_proposta(self):
        """Gera Proposta (.pptx) e em seguida converte em PDF."""
        abas_ativas = self.get_abas_ativas()
        if not abas_ativas:
            showerror("Erro", "Nenhuma aba criada para gerar Proposta.")
            return
        pptx_file = gerar_proposta(
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        if pptx_file and os.path.exists(pptx_file):
            upload_pptx_and_export_to_pdf(pptx_file)

    def on_gerar_mat_tecnico(self):
        """Gera Material Técnico (.pptx) e em seguida converte em PDF."""
        abas_ativas = self.get_abas_ativas()
        if not abas_ativas:
            showerror("Erro", "Nenhuma aba criada para gerar Material Técnico.")
            return
        pptx_file = gerar_material(
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        if pptx_file and os.path.exists(pptx_file):
            upload_pptx_and_export_to_pdf(pptx_file)

    def on_gerar_tudo(self):
        """Gera Proposta e Material Técnico, cada um em .pptx, depois converte em PDF."""
        abas_ativas = self.get_abas_ativas()
        if not abas_ativas:
            showerror("Erro", "Nenhuma aba criada para gerar.")
            return

        # 1) Gera Proposta
        pptx_prop = gerar_proposta(
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        if pptx_prop and os.path.exists(pptx_prop):
            upload_pptx_and_export_to_pdf(pptx_prop)

        # 2) Gera Material
        pptx_mat = gerar_material(
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        if pptx_mat and os.path.exists(pptx_mat):
            upload_pptx_and_export_to_pdf(pptx_mat)


def main():
    app = MainApp()
    app.mainloop()

if __name__ == "__main__":
    main()
