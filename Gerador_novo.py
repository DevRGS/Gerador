# coding: utf-8
import sys
import subprocess
import importlib

# Dicionário: "nome para importar" : "nome do pacote no PyPI"
dependencias = {
    "ttkbootstrap": "ttkbootstrap",
    "pptx": "python-pptx",
    "googleapiclient": "google-api-python-client",
    "google.auth.transport": "google-auth-httplib2",
    "google_auth_oauthlib": "google-auth-oauthlib",
    "requests": "requests"
}

for modulo, pacote in dependencias.items():
    try:
        importlib.import_module(modulo)
    except ImportError:
        print(f"Instalando '{pacote}' pois o módulo '{modulo}' não foi encontrado...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", pacote])
        # Atualiza o cache de importação
        importlib.invalidate_caches()
        try:
            importlib.import_module(modulo)
        except ImportError:
            print(f"Falha ao importar '{modulo}' mesmo após a instalação.")

import os
import requests
import io
import pickle
import json
from datetime import date

# Tkinter e ttkbootstrap
import tkinter as tk
from tkinter import ttk
import ttkbootstrap as ttkb
from tkinter.messagebox import showerror, showinfo

# python-pptx
from pptx import Presentation

# Google Drive / Auth
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload

def baixar_arquivo_if_needed(nome_arquivo, url):
    if not os.path.exists(nome_arquivo):
        print(f"Baixando {nome_arquivo}...")
        try:
            r = requests.get(url, timeout=30) # Timeout aumentado
            r.raise_for_status() # Verifica erros HTTP
            with open(nome_arquivo, "wb") as f:
                f.write(r.content)
            print(f"{nome_arquivo} baixado com sucesso.")
        except requests.exceptions.RequestException as e:
             showerror("Erro de Download", f"Não foi possível baixar '{nome_arquivo}'.\nVerifique sua conexão ou a URL.\nErro: {e}")
             # Decide se quer parar a execução ou apenas avisar
             # raise e # Descomente para parar se o download for crítico


# ---------------------------------------------------------
# Ajustes Globais
# ---------------------------------------------------------
# Tenta encontrar o diretório do script de forma mais robusta
if getattr(sys, 'frozen', False):
    # Se rodando como executável (pyinstaller)
    script_dir = os.path.dirname(sys.executable)
elif __file__:
    # Se rodando como script .py
    script_dir = os.path.dirname(os.path.abspath(__file__))
else:
    # Fallback para diretório atual
    script_dir = os.getcwd()

try:
    os.chdir(script_dir)
    print("Diretório de trabalho definido para:", os.getcwd())
except Exception as e:
    print(f"Aviso: Não foi possível mudar para o diretório do script: {e}")
    print(f"Diretório atual: {os.getcwd()}")


CONFIG_FILE = "config_vendedor.json"
MAX_ABAS = 10

# ---------------------------------------------------------
# Configurações de vendedor (salvar/carregar)
# ---------------------------------------------------------
def carregar_config(nome_closer_var, celular_closer_var, email_closer_var):
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            nome_closer_var.set(data.get("nome_vendedor", ""))
            celular_closer_var.set(data.get("celular_vendedor", ""))
            email_closer_var.set(data.get("email_vendedor", ""))
        except json.JSONDecodeError:
            print(f"Aviso: Arquivo '{CONFIG_FILE}' parece corrompido. Ignorando.")
            pass # Ignora arquivo corrompido
        except Exception as e:
            print(f"Erro desconhecido ao carregar config: {e}")

def salvar_config(nome_closer, celular_closer, email_closer):
    dados = {
        "nome_vendedor": nome_closer,
        "celular_vendedor": celular_closer,
        "email_vendedor": email_closer
    }
    # Tenta tornar o arquivo gravável se existir (útil em alguns sistemas)
    if os.path.exists(CONFIG_FILE):
        try:
            os.chmod(CONFIG_FILE, 0o666)
        except OSError:
            pass # Ignora erro se não puder mudar permissão
    try:
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(dados, f, indent=4, ensure_ascii=False)
    except PermissionError:
        print(f"Aviso: Sem permissão para salvar {CONFIG_FILE}")
    except Exception as e:
        print(f"Erro ao salvar config: {e}")


# ---------------------------------------------------------
# Dados de Planos e Tabelas de Preço - Base Mensal = Custo Efetivo no Anual
# ---------------------------------------------------------
LISTA_PLANOS_UI = ["PDV Básico", "Gestão", "Performance", "Autoatendimento", "Bling", "Em Branco"]
LISTA_PLANOS_BLING = ["Bling - Básico", "Bling - Com Estoque em Grade"]

PLAN_INFO = {
    "PDV Básico": {
        # base_mensal AGORA significa -> Custo Mensal Efetivo no contrato Anual
        "base_mensal": 110.00,
        "min_pdv": 1, "min_users": 2,
        "max_extra_users": 1, "max_extra_pdvs": 0,
        "mandatory": ["Usuários", "30 Notas Fiscais", "Suporte Técnico - Via chamados", "Relatório Básico", "PDV - Frente de Caixa"],
        "allowed_optionals": ["Smart Menu", "Terminais Autoatendimento", "Hub de Delivery", "Delivery Direto Profissional", "Delivery Direto VIP", "TEF", "Importação de XML", "Cardápio digital"]
    },
    "Gestão": {
        "base_mensal": 221.11,
        "min_pdv": 2, "min_users": 3,
        "max_extra_users": 2, "max_extra_pdvs": 1,
        "mandatory": ["Notas Fiscais Ilimitadas", "Importação de XML", "PDV - Frente de Caixa", "Usuários", "Painel Senha TV", "Estoque em Grade", "Relatórios", "Suporte Técnico - Via chamados", "Suporte Técnico - Via chat", "Delivery", "Relatório KDS"],
        "allowed_optionals": ["Facilita NFE", "Conciliação Bancária", "Contratos de cartões e outros", "Delivery Direto Profissional", "Delivery Direto VIP", "TEF", "Integração API", "Business Intelligence (BI)", "Backup Realtime", "Cardápio digital", "Smart Menu", "Hub de Delivery", "Ordem de Serviço", "App Gestão CPlug", "Painel Senha Mobile", "Controle de Mesas", "Produção", "Promoções", "Marketing", "Relatório Dinâmico", "Atualização em tempo real", "Smart TEF", "Terminais Autoatendimento", "Suporte Técnico - Estendido"]
    },
    "Performance": {
        "base_mensal": 554.44,
        "min_pdv": 3, "min_users": 5,
        "max_extra_users": 5, "max_extra_pdvs": 2,
        "mandatory": ["Produção", "Promoções", "Notas Fiscais Ilimitadas", "Importação de XML", "Hub de Delivery", "Ordem de Serviço", "Delivery", "App Gestão CPlug", "Relatório KDS", "Painel Senha TV", "Painel Senha Mobile", "Controle de Mesas", "Estoque em Grade", "Marketing", "Relatórios", "Relatório Dinâmico", "Atualização em tempo real", "Facilita NFE", "Conciliação Bancária", "Contratos de cartões e outros", "Suporte Técnico - Via chamados", "Suporte Técnico - Via chat", "Suporte Técnico - Estendido", "PDV - Frente de Caixa", "Smart TEF", "Usuários"],
        "allowed_optionals": ["TEF", "Programa de Fidelidade", "Integração Tap", "Integração API", "Business Intelligence (BI)", "Backup Realtime", "Cardápio digital", "Smart Menu", "Terminais Autoatendimento", "Delivery Direto Profissional", "Delivery Direto VIP"]
    },
    "Autoatendimento": {
        "base_mensal": 419.90, # Assumindo que este é o Custo Mensal Efetivo no anual para 1 terminal
        "min_pdv": 0, "min_users": 1,
        "max_extra_users": 998, "max_extra_pdvs": 99,
        "mandatory": ["Contratos de cartões e outros", "Estoque em Grade", "Notas Fiscais Ilimitadas", "Produção"], # Vendas.. removido
        "allowed_optionals": [] # Definir se houver
    },
    "Bling - Básico": {
        "base_mensal": 189.90, # Este JÁ É o valor mensal efetivo no anual
        "min_pdv": 1, "min_users": 5,
        "max_extra_users": 994, "max_extra_pdvs": 98,
        "mandatory": ["Relatórios", "Notas Fiscais Ilimitadas"], # Vendas.. removido
        "allowed_optionals": [], # Definir se houver
        # "base_mensal_original": 369.80, # Pode ser útil para exibir "De R$..."
    },
    "Bling - Com Estoque em Grade": {
        "base_mensal": 219.90, # Este JÁ É o valor mensal efetivo no anual
        "min_pdv": 1, "min_users": 5,
        "max_extra_users": 994, "max_extra_pdvs": 98,
        "mandatory": ["Relatórios", "Notas Fiscais Ilimitadas", "Estoque em Grade"], # Vendas.. removido
        "allowed_optionals": [], # Definir se houver
        # "base_mensal_original": 399.80, # Pode ser útil para exibir "De R$..."
    },
    "Em Branco": {
        "base_mensal": 0.0,
        "min_pdv": 0, "min_users": 0,
        "max_extra_users": 999, "max_extra_pdvs": 99,
        "mandatory": [],
        "allowed_optionals": list(precos_mensais.keys()) + ["PDV - Frente de Caixa", "Usuários", "Smart TEF"] # Permite tudo?
    }
}

# Módulos SEM DESCONTO (se a lógica de desconto manual ainda for relevante)
# Itens cujo preço mensal não deve ser afetado pelo % de desconto editado
SEM_DESCONTO = {
    "TEF", "Terminais Autoatendimento", "Smart TEF",
    "Delivery Direto Profissional", "Delivery Direto VIP",
    "Programa de Fidelidade", "Integração Tap", "Integração API",
    "Business Intelligence (BI)", "Backup Realtime",
    # Manter os antigos se ainda relevantes?
    # "Domínio Próprio", ... etc
}

# Dicionário de preços mensais dos EXTRAS (mantido da versão anterior)
precos_mensais = {
    # Módulos que são *sempre* fixos não precisam de preço aqui
    # Preços para módulos que podem ser OPCIONAIS em alguns planos
    "Importação de XML": 29.00,
    "Produção": 30.00,
    "Promoções": 24.50,
    "Hub de Delivery": 79.00,
    "Ordem de Serviço": 20.00,
    "App Gestão CPlug": 20.00,
    "Painel Senha Mobile": 49.00,
    "Controle de Mesas": 49.00,
    "Marketing": 24.50,
    "Relatório Dinâmico": 50.00,
    "Atualização em tempo real": 49.00,
    "Facilita NFE": 99.00,
    "Conciliação Bancária": 50.00,
    "Contratos de cartões e outros": 50.00,
    "Suporte Técnico - Estendido": 99.00,
    "Smart TEF": 49.90, # Usado para cálculo de extras no Gestão e Performance

    # Opcionais puros
    "Smart Menu": 99.90,
    "Terminais Autoatendimento": 199.00, # Custo por terminal extra
    "Delivery Direto Profissional": 200.00,
    "Delivery Direto VIP": 300.00,
    "TEF": 99.90,
    "Cardápio digital": 99.00,
    "Backup Realtime": 199.90,
    "Business Intelligence (BI)": 199.00,
    "Programa de Fidelidade": 299.90,
    "Integração Tap": 299.00,
    "Integração API": 299.00,

    # Itens talvez obsoletos ou implícitos em outros
    # "60 Notas Fiscais": 40.00, ... etc
    # Manter apenas os que podem ser selecionados explicitamente
}

# Custos Adicionais (fixos por item extra)
PRECO_EXTRA_USUARIO = 19.00
PRECO_EXTRA_PDV_GESTAO_PERFORMANCE = 59.90
PRECO_EXTRA_PDV_BLING = 40.00 # Se Bling ainda permitir extra PDV

# ---------------------------------------------------------
# *** CORRIGIDO *** Função utilitária para substituir placeholders no Slide
# ---------------------------------------------------------
def substituir_placeholders_no_slide(slide, dados):
    """
    Substitui os placeholders (chaves do dicionário 'dados') pelo seu valor
    correspondente em todos os text frames do slide.
    Preserva a formatação original do run onde o placeholder está.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            # É crucial iterar sobre os runs originais
            inline_runs_text = "" # Para reconstruir o texto do parágrafo se necessário (raro)
            for run in paragraph.runs:
                # Preserva o texto original do run para comparação
                original_text = run.text
                modified_text = original_text # Começa com o texto original

                # Tenta substituir cada placeholder no texto deste run
                for k, v in dados.items():
                    placeholder = f"{{{k}}}" # Assume placeholders como {chave}
                    # Garante que v seja string
                    v_str = str(v) if v is not None else ""
                    # Faz a substituição no texto modificado
                    # Usar replace() diretamente no run.text pode ser problemático se
                    # um placeholder estiver dividido entre runs, mas é o mais comum
                    if placeholder in modified_text:
                        modified_text = modified_text.replace(placeholder, v_str)

                # Apenas atualiza o run.text se houve mudança
                if modified_text != original_text:
                    run.text = modified_text

                # Adiciona ao texto reconstruído (apenas para debug ou casos complexos)
                # inline_runs_text += run.text

            # Opcional: Verificar se o texto completo do parágrafo corresponde
            # if shape.text_frame.text != inline_runs_text:
            #     print(f"Aviso: Discrepância de texto no parágrafo após substituição. Shape: {shape.name}")

# ---------------------------------------------------------
# Classe PlanoFrame (Aba)
# ---------------------------------------------------------
class PlanoFrame(ttkb.Frame):
    def __init__(
        self,
        parent,
        aba_index,
        nome_cliente_var_shared,
        validade_proposta_var_shared,
        on_close_callback=None
    ):
        super().__init__(parent)
        self.aba_index = aba_index
        self.on_close_callback = on_close_callback

        # Variáveis compartilhadas
        self.nome_cliente_var = nome_cliente_var_shared
        self.validade_proposta_var = validade_proposta_var_shared
        self.nome_plano_var = tk.StringVar(value="") # Nome editável do plano

        self.current_plan = "PDV Básico" # Plano default
        self.spin_pdv_var = tk.IntVar(value=1)
        self.spin_users_var = tk.IntVar(value=1)
        self.spin_smart_tef_var = tk.IntVar(value=0)
        self.spin_terminais_aa_var = tk.IntVar(value=0) # Novo spinbox para Autoatendimento

        # Módulos (checkboxes)
        self.modules = {
            # Fixos (necessários para lógica, mas controlados por código/plano)
            "Usuários": tk.IntVar(), "30 Notas Fiscais": tk.IntVar(), "Suporte Técnico - Via chamados": tk.IntVar(),
            "Relatório Básico": tk.IntVar(), "PDV - Frente de Caixa": tk.IntVar(), "Notas Fiscais Ilimitadas": tk.IntVar(),
            "Importação de XML": tk.IntVar(), "Painel Senha TV": tk.IntVar(), "Estoque em Grade": tk.IntVar(),
            "Relatórios": tk.IntVar(), "Suporte Técnico - Via chat": tk.IntVar(), "Delivery": tk.IntVar(),
            "Relatório KDS": tk.IntVar(), "Produção": tk.IntVar(), "Promoções": tk.IntVar(), "Hub de Delivery": tk.IntVar(),
            "Ordem de Serviço": tk.IntVar(), "App Gestão CPlug": tk.IntVar(), "Painel Senha Mobile": tk.IntVar(),
            "Controle de Mesas": tk.IntVar(), "Marketing": tk.IntVar(), "Relatório Dinâmico": tk.IntVar(),
            "Atualização em tempo real": tk.IntVar(), "Facilita NFE": tk.IntVar(), "Conciliação Bancária": tk.IntVar(),
            "Contratos de cartões e outros": tk.IntVar(), "Suporte Técnico - Estendido": tk.IntVar(), "Smart TEF": tk.IntVar(),
            # Opcionais (reais checkboxes que aparecem na UI)
            "Smart Menu": tk.IntVar(),
            # "Terminais Autoatendimento": tk.IntVar(), # Removido - Controlado pelo spinbox
            "Delivery Direto Profissional": tk.IntVar(),
            "Delivery Direto VIP": tk.IntVar(), "TEF": tk.IntVar(), "Cardápio digital": tk.IntVar(),
            "Integração API": tk.IntVar(), "Business Intelligence (BI)": tk.IntVar(), "Backup Realtime": tk.IntVar(),
            "Programa de Fidelidade": tk.IntVar(), "Integração Tap": tk.IntVar(),
        }
        self.check_buttons = {} # Dicionário para guardar os widgets Checkbutton

        # Overrides de cálculo
        self.user_override_anual_active = tk.BooleanVar(value=False)
        self.user_override_discount_active = tk.BooleanVar(value=False)
        self.valor_anual_editavel = tk.StringVar(value="0.00") # Representa o TOTAL ANUAL PAGO ADIANTADO
        self.desconto_personalizado = tk.StringVar(value="10") # Começa com 10% padrão

        # Armazenar valores calculados
        self.computed_mensal_sem_fidelidade = 0.0
        self.computed_mensal_efetivo_anual = 0.0
        self.computed_anual_total = 0.0
        self.computed_desconto_percent = 0.0
        self.computed_custo_adicional = 0.0

        # --- Layout com Scrollbar ---
        self.canvas = tk.Canvas(self)
        self.canvas.pack(side="left", fill="both", expand=True)
        self.scrollbar = ttkb.Scrollbar(self, orient="vertical", command=self.canvas.yview)
        self.scrollbar.pack(side="right", fill="y")
        self.canvas.configure(yscrollcommand=self.scrollbar.set)
        self.container = ttkb.Frame(self.canvas)
        self.canvas_window = self.canvas.create_window((0,0), window=self.container, anchor="nw") # Guardar ref
        self.container.bind("<Configure>", lambda e: self.canvas.configure(scrollregion=self.canvas.bbox("all")))
        # Bind para mouse wheel funcionar no canvas
        self.canvas.bind_all("<MouseWheel>", lambda e: self.canvas.yview_scroll(int(-1*(e.delta/120)), "units"))
        # Bind para o container também, caso o foco esteja nele
        self.container.bind("<MouseWheel>", lambda e: self.canvas.yview_scroll(int(-1*(e.delta/120)), "units"))

        self.frame_main = ttkb.Frame(self.container)
        self.frame_main.pack(fill="both", expand=True)
        self.frame_left = ttkb.Frame(self.frame_main)
        self.frame_left.pack(side="left", fill="both", expand=True, padx=5, pady=5)
        self.frame_right = ttkb.Frame(self.frame_main)
        self.frame_right.pack(side="left", fill="y", padx=5, pady=5)
        # --- Fim Layout Scrollbar ---

        self._montar_layout_esquerda()
        self._montar_layout_direita()
        self.configurar_plano("PDV Básico") # Configura o plano inicial

        # Ajustar tamanho inicial do canvas
        self.container.update_idletasks()
        self.canvas.config(scrollregion=self.canvas.bbox("all"))


    def fechar_aba(self):
        if self.on_close_callback:
            self.on_close_callback(self.aba_index)

    def on_bling_selected(self, event=None):
        selected_bling_plan = self.bling_var.get()
        if selected_bling_plan in LISTA_PLANOS_BLING:
            self.configurar_plano(selected_bling_plan)
        # Não reseta o combobox para manter a seleção visível
        # self.bling_var.set("Selecionar Bling...")

    def _montar_layout_esquerda(self):
        # Top Bar e Seleção de Planos
        top_bar = ttkb.Frame(self.frame_left)
        top_bar.pack(fill="x", pady=5)
        ttkb.Label(top_bar, text=f"Aba Plano {self.aba_index}", font="-size 12 -weight bold").pack(side="left")
        btn_close = ttkb.Button(top_bar, text="X", command=self.fechar_aba, bootstyle="danger-outline", width=3)
        btn_close.pack(side="right")

        frame_planos = ttkb.Labelframe(self.frame_left, text="Selecionar Plano Base")
        frame_planos.pack(fill="x", pady=5)
        self.bling_combobox = None
        plan_buttons_frame = ttkb.Frame(frame_planos) # Frame para botões
        plan_buttons_frame.pack(fill="x")
        for i, p in enumerate(LISTA_PLANOS_UI):
            if p == "Bling":
                self.bling_var = tk.StringVar(value="Selecionar Bling...")
                self.bling_combobox = ttk.Combobox(plan_buttons_frame, textvariable=self.bling_var,
                                                   values=LISTA_PLANOS_BLING, state="readonly", width=25)
                self.bling_combobox.grid(row=0, column=i, padx=5, pady=5, sticky="ew")
                self.bling_combobox.bind("<<ComboboxSelected>>", self.on_bling_selected)
            else:
                btn = ttkb.Button(plan_buttons_frame, text=p, width=15,
                                  command=lambda pl=p: self.configurar_plano(pl))
                btn.grid(row=0, column=i, padx=5, pady=5, sticky="ew") # Usa grid para melhor alinhamento
            plan_buttons_frame.grid_columnconfigure(i, weight=1) # Faz colunas expansíveis


        # Módulos Opcionais (Checkboxes)
        frame_mod = ttkb.Labelframe(self.frame_left, text="Módulos Opcionais (Marque para adicionar)")
        frame_mod.pack(fill="both", expand=True, pady=5)
        f_mod_cols = ttkb.Frame(frame_mod)
        f_mod_cols.pack(fill="both", expand=True)

        f_mod_left = ttkb.Frame(f_mod_cols)
        f_mod_left.pack(side="left", fill="both", expand=True, padx=5)
        f_mod_right = ttkb.Frame(f_mod_cols)
        f_mod_right.pack(side="left", fill="both", expand=True, padx=5)

        # Lista de módulos que podem aparecer como checkbox na UI
        # (Todos os que não são controlados por spinbox e não são implicitamente mandatórios em TODOS os planos)
        displayable_mods_ui = sorted([
            m for m, var in self.modules.items() if m not in [
                "PDV - Frente de Caixa", "Usuários", "Smart TEF", # Controlados por spinbox
                "Terminais Autoatendimento", # Controlado por spinbox
                "Relatórios", "Relatório Básico", # Mandatórios implícitos ou genéricos
                "30 Notas Fiscais", "Notas Fiscais Ilimitadas", # Ligados a planos
                "Suporte Técnico - Via chamados", # Mandatório base
            ]
        ])

        mid = len(displayable_mods_ui)//2
        left_side = displayable_mods_ui[:mid]
        right_side = displayable_mods_ui[mid:]
        self.check_buttons = {} # Limpa para garantir que só os criados agora estarão lá

        for m in left_side:
             if m in self.modules: # Confirma que o módulo existe no dicionário principal
                 cb = ttk.Checkbutton(f_mod_left, text=m, variable=self.modules[m], command=self.atualizar_valores)
                 cb.pack(anchor="w", pady=2)
                 self.check_buttons[m] = cb # Armazena o widget

        for m in right_side:
             if m in self.modules:
                 cb = ttk.Checkbutton(f_mod_right, text=m, variable=self.modules[m], command=self.atualizar_valores)
                 cb.pack(anchor="w", pady=2)
                 self.check_buttons[m] = cb

        # Frame Dados Cliente e Plano
        frame_dados = ttkb.Labelframe(self.frame_left, text="Dados do Cliente e Proposta")
        frame_dados.pack(fill="x", pady=5)
        ttkb.Label(frame_dados, text="Nome do Cliente:").grid(row=0, column=0, sticky="w", padx=5, pady=2)
        ttkb.Entry(frame_dados, textvariable=self.nome_cliente_var, width=30).grid(row=0, column=1, padx=5, pady=2)
        ttkb.Label(frame_dados, text="Validade Proposta:").grid(row=1, column=0, sticky="w", padx=5, pady=2)
        ttkb.Entry(frame_dados, textvariable=self.validade_proposta_var, width=15).grid(row=1, column=1, padx=5, pady=2, sticky="w")
        ttkb.Label(frame_dados, text="Nome do Plano (Opcional):").grid(row=2, column=0, sticky="w", padx=5, pady=2)
        ttkb.Entry(frame_dados, textvariable=self.nome_plano_var, width=30).grid(row=2, column=1, padx=5, pady=2)


    def _montar_layout_direita(self):
        # Quantidades (PDV, Usuários, Smart TEF, Terminais AA)
        frame_inc = ttkb.Labelframe(self.frame_right, text="Quantidades")
        frame_inc.pack(fill="x", pady=5)

        ttkb.Label(frame_inc, text="PDVs - Frente de Caixa").grid(row=0, column=0, sticky="w", padx=5, pady=2)
        self.sp_pdv = ttkb.Spinbox(frame_inc, from_=0, to=99, textvariable=self.spin_pdv_var, width=5, command=self.atualizar_valores)
        self.sp_pdv.grid(row=0, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="Usuários").grid(row=1, column=0, sticky="w", padx=5, pady=2)
        self.sp_usr = ttkb.Spinbox(frame_inc, from_=0, to=999, textvariable=self.spin_users_var, width=5, command=self.atualizar_valores)
        self.sp_usr.grid(row=1, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="Smart TEF").grid(row=2, column=0, sticky="w", padx=5, pady=2)
        self.sp_smf = ttkb.Spinbox(frame_inc, from_=0, to=99, textvariable=self.spin_smart_tef_var, width=5, command=self.atualizar_valores)
        self.sp_smf.grid(row=2, column=1, padx=5, pady=2)

        ttkb.Label(frame_inc, text="Terminais Autoatendimento").grid(row=3, column=0, sticky="w", padx=5, pady=2)
        self.sp_taa = ttkb.Spinbox(frame_inc, from_=0, to=99, textvariable=self.spin_terminais_aa_var, width=5, command=self.atualizar_valores)
        self.sp_taa.grid(row=3, column=1, padx=5, pady=2)


        # --- Frame Valores Finais (Layout Atualizado) ---
        frame_valores = ttkb.Labelframe(self.frame_right, text="Valores da Proposta")
        frame_valores.pack(fill="both", pady=5, expand=True) # Expandir para preencher espaço

        # 1. Mensal (Sem Fidelidade)
        self.lbl_plano_mensal_sem_fid = ttkb.Label(frame_valores, text="Mensal (Sem Fidelidade): R$ 0,00", font="-size 11")
        self.lbl_plano_mensal_sem_fid.pack(pady=(5, 2), anchor="w", padx=5)

        # 2. Custo Treinamento/Implementação (Associado ao mensal sem fidelidade)
        self.lbl_treinamento = ttkb.Label(frame_valores, text="+ Custo Treinamento: R$ 0,00", font="-size 9")
        self.lbl_treinamento.pack(pady=(0, 5), anchor="w", padx=15) # Indentado

        # 3. Mensal (No Plano Anual)
        self.lbl_plano_mensal_no_anual = ttkb.Label(frame_valores, text="Mensal (no Plano Anual): R$ 0,00", font="-size 12 -weight bold")
        self.lbl_plano_mensal_no_anual.pack(pady=5, anchor="w", padx=5)

        # 4. Anual (Pagamento Único)
        self.lbl_plano_anual_total = ttkb.Label(frame_valores, text="Anual (Pagamento Único): R$ 0,00", font="-size 12 -weight bold")
        self.lbl_plano_anual_total.pack(pady=5, anchor="w", padx=5)

        # 5. Desconto Aplicado (Informativo)
        self.lbl_desconto = ttkb.Label(frame_valores, text="Desconto Anual Aplicado: 10%", font="-size 9")
        self.lbl_desconto.pack(pady=(5, 10), anchor="w", padx=5)

        # Separador
        ttk.Separator(frame_valores, orient='horizontal').pack(fill='x', pady=5, padx=5)

        # --- Frames de Edição Manual ---
        frame_edicao = ttkb.Frame(frame_valores)
        frame_edicao.pack(fill="x", pady=5)

        # 6. Edição Anual Total
        frame_edit_anual = ttkb.Labelframe(frame_edicao, text="Editar Anual Total (R$)")
        frame_edit_anual.pack(side="left", padx=5, fill="x", expand=True)
        e_anual = ttkb.Entry(frame_edit_anual, textvariable=self.valor_anual_editavel, width=10)
        e_anual.pack(side="left", padx=5, pady=2)
        e_anual.bind("<KeyRelease>", self.on_user_edit_valor_anual)
        b_reset_anual = ttkb.Button(frame_edit_anual, text="Reset", command=self.on_reset_anual, width=5, bootstyle="warning-outline")
        b_reset_anual.pack(side="left", padx=5, pady=2)

        # 7. Edição Desconto %
        frame_edit_desc = ttkb.Labelframe(frame_edicao, text="Editar Desconto (%)")
        frame_edit_desc.pack(side="left", padx=5, fill="x", expand=True)
        e_desc = ttkb.Entry(frame_edit_desc, textvariable=self.desconto_personalizado, width=5)
        e_desc.pack(side="left", padx=5, pady=2)
        e_desc.bind("<KeyRelease>", self.on_user_edit_desconto)
        b_reset_desc = ttkb.Button(frame_edit_desc, text="Reset", command=self.on_reset_desconto, width=5, bootstyle="warning-outline")
        b_reset_desc.pack(side="left", padx=5, pady=2)


    # --- Funções de Edição/Reset ---
    def on_user_edit_valor_anual(self, *args):
        # Usuário editou o VALOR TOTAL ANUAL
        self.user_override_anual_active.set(True)
        self.user_override_discount_active.set(False) # Desativa override de desconto
        self.atualizar_valores() # Recalcula tudo e atualiza UI

    def on_reset_anual(self):
        self.user_override_anual_active.set(False)
        # Não reseta valor_anual_editavel aqui, deixa atualizar_valores recalcular
        self.atualizar_valores()

    def on_user_edit_desconto(self, *args):
        # Usuário editou o PERCENTUAL DE DESCONTO
        self.user_override_discount_active.set(True)
        self.user_override_anual_active.set(False) # Desativa override de valor anual
        self.atualizar_valores() # Recalcula tudo e atualiza UI

    def on_reset_desconto(self):
        self.user_override_discount_active.set(False)
        # Não reseta desconto_personalizado aqui, deixa atualizar_valores recalcular (para 10% ou calculado)
        self.atualizar_valores()

    def configurar_plano(self, plano):
        # Reset Bling Combobox se outro plano for selecionado
        if not plano.startswith("Bling -") and self.bling_combobox:
             self.bling_var.set("Selecionar Bling...")

        if plano not in PLAN_INFO:
             showerror("Erro de Configuração", f"Plano '{plano}' não encontrado nas definições.")
             return

        info = PLAN_INFO[plano]
        self.current_plan = plano

        # --- Configurar Spinboxes (Mínimos e Máximos) ---
        min_pdv = info.get("min_pdv", 0); max_pdv = 99 # Simplificado, poderia ter max_extra
        min_users = info.get("min_users", 0); max_users = 999

        self.spin_pdv_var.set(min_pdv)
        self.sp_pdv.config(from_=min_pdv, to=max_pdv)
        self.spin_users_var.set(min_users)
        self.sp_usr.config(from_=min_users, to=max_users)

        # Limite Smart TEF no Gestão (e incluído no Performance)
        min_smart_tef = 0; max_smart_tef = 99 # Padrão
        val_inicial_smart_tef = 0
        if plano == "Gestão":
             max_smart_tef = 3 # Permite adicionar até 3
        elif plano == "Performance":
             min_smart_tef = 3 # Performance já inclui 3
             max_smart_tef = 3 # Não permite adicionar mais pelo spinbox
             val_inicial_smart_tef = 3

        self.spin_smart_tef_var.set(val_inicial_smart_tef)
        self.sp_smf.config(from_=min_smart_tef, to=max_smart_tef)
        self.sp_smf.config(state='normal' if max_smart_tef > min_smart_tef else 'disabled') # Desabilita se fixo

        # Configurar Terminais de Autoatendimento
        min_taa = 0; max_taa = 99
        val_inicial_taa = 0
        if plano == "Autoatendimento":
            min_taa = 1 # Plano Autoatendimento começa com 1
            val_inicial_taa = 1
        self.spin_terminais_aa_var.set(val_inicial_taa)
        self.sp_taa.config(from_=min_taa, to=max_taa)
        # Habilita/Desabilita TAA: Permitido no PDV Básico, Gestão, Performance e Autoatendimento
        # (Talvez permitir em Em Branco também?)
        allow_taa = plano in ["PDV Básico", "Gestão", "Performance", "Autoatendimento", "Em Branco"]
        self.sp_taa.config(state='normal' if allow_taa else 'disabled')


        # --- Resetar/Configurar Módulos (Checkboxes) ---
        # 1. Limpa todos os checkboxes e reseta estado
        for m, var in self.modules.items():
            var.set(0)
            if m in self.check_buttons:
                self.check_buttons[m].config(state='normal') # Habilita todos por padrão

        # 2. Define os mandatórios e desabilita seus checkboxes
        mandatory = info.get("mandatory", [])
        for m in mandatory:
            if m in self.modules:
                self.modules[m].set(1) # Marca como selecionado
                if m in self.check_buttons:
                    self.check_buttons[m].config(state='disabled') # Desabilita o checkbox

        # 3. Desabilita checkboxes que NÃO são opcionais permitidos para este plano
        allowed = info.get("allowed_optionals", [])
        # Considera também os mandatórios como "permitidos" (mas já desabilitados)
        implicitly_allowed = set(mandatory)
        # Itens controlados por Spinbox não devem ser controlados por Checkbox
        spinbox_controlled = {"PDV - Frente de Caixa", "Usuários", "Smart TEF", "Terminais Autoatendimento"}

        for m, cb in self.check_buttons.items():
            if m not in allowed and m not in implicitly_allowed:
                 # Se não está na lista de opcionais permitidos NEM é mandatório
                 cb.config(state='disabled')
                 self.modules[m].set(0) # Garante que não fique marcado se não for permitido
            elif m in spinbox_controlled:
                 # Garante que itens de spinbox não sejam clicáveis como checkbox
                 cb.config(state='disabled')
                 # O valor (set(1)) é controlado pela lógica do spinbox ou mandatório

        # --- Resetar overrides e recalcular ---
        self.user_override_anual_active.set(False)
        self.user_override_discount_active.set(False)
        # Definir o desconto padrão como 10 ao configurar o plano
        self.desconto_personalizado.set("10")
        self.atualizar_valores() # Dispara o recálculo inicial

    def _calcular_extras(self):
        """Calcula o custo MENSAL total dos extras e separa por descontável/não descontável."""
        total_extras_cost = 0.0
        total_extras_descontavel = 0.0
        total_extras_nao_descontavel = 0.0
        info = PLAN_INFO.get(self.current_plan, {}) # Pega info do plano atual
        if not info: return 0.0, 0.0, 0.0 # Retorna zero se plano inválido

        mandatory = info.get("mandatory", [])
        base_mensal_efetivo = info.get("base_mensal", 0.0) # Custo base efetivo anual do plano

        # --- 1. Itens controlados por Spinbox ---

        # PDVs Extras
        pdv_atuais = self.spin_pdv_var.get()
        pdv_incluidos = info.get("min_pdv", 0)
        pdv_extras = max(0, pdv_atuais - pdv_incluidos)
        if pdv_extras > 0:
            pdv_price = 0.0
            if self.current_plan in ["Gestão", "Performance"]: pdv_price = PRECO_EXTRA_PDV_GESTAO_PERFORMANCE
            elif self.current_plan.startswith("Bling"): pdv_price = PRECO_EXTRA_PDV_BLING
            # PDV Básico não permite PDV extra (max_extra_pdvs: 0)
            # Plano "Em Branco" - Qual preço usar? Usar o de Gestão/Performance por enquanto.
            elif self.current_plan == "Em Branco": pdv_price = PRECO_EXTRA_PDV_GESTAO_PERFORMANCE

            cost_pdv_extra = pdv_extras * pdv_price
            total_extras_cost += cost_pdv_extra
            # Assumindo que PDV extra é descontável, exceto se explicitamente listado em SEM_DESCONTO
            if "PDV Extra" not in SEM_DESCONTO: total_extras_descontavel += cost_pdv_extra
            else: total_extras_nao_descontavel += cost_pdv_extra

        # Users Extras
        users_atuais = self.spin_users_var.get()
        users_incluidos = info.get("min_users", 0)
        users_extras = max(0, users_atuais - users_incluidos)
        if users_extras > 0:
            cost_users_extra = users_extras * PRECO_EXTRA_USUARIO
            total_extras_cost += cost_users_extra
            if "User Extra" not in SEM_DESCONTO: total_extras_descontavel += cost_users_extra
            else: total_extras_nao_descontavel += cost_users_extra

        # Smart TEF Extras (Apenas no Gestão, pois Performance já inclui 3)
        if self.current_plan == "Gestão":
            smart_tef_atuais = self.spin_smart_tef_var.get()
            smart_tef_incluidos = 0 # Gestão não inclui nenhum fixo
            smart_tef_extras = max(0, smart_tef_atuais - smart_tef_incluidos)
            if smart_tef_extras > 0:
                price = smart_tef_extras * precos_mensais.get("Smart TEF", 0.0)
                total_extras_cost += price
                if "Smart TEF" not in SEM_DESCONTO: total_extras_descontavel += price
                else: total_extras_nao_descontavel += price
        # No Performance, os 3 já estão no custo base, não adiciona aqui.

        # Terminais Autoatendimento Extras
        taa_atuais = self.spin_terminais_aa_var.get()
        taa_incluidos = 1 if self.current_plan == "Autoatendimento" else 0
        taa_extras = max(0, taa_atuais - taa_incluidos)
        if taa_extras > 0:
            # Preço por terminal extra
            price_per_taa_extra = precos_mensais.get("Terminais Autoatendimento", 199.00)
            cost_taa_extra = taa_extras * price_per_taa_extra
            total_extras_cost += cost_taa_extra
            # Terminais AA estão em SEM_DESCONTO
            total_extras_nao_descontavel += cost_taa_extra

        # --- 2. Módulos Extras (Checkboxes Opcionais Selecionados) ---
        for m, var_m in self.modules.items():
             # Considera apenas os checkboxes VISÍVEIS na UI e que estão MARCADOS
             if m in self.check_buttons and var_m.get() == 1:
                 # Verifica se NÃO é um módulo mandatório para ESTE plano
                 if m not in mandatory:
                     price = precos_mensais.get(m, 0.0)
                     total_extras_cost += price
                     if m not in SEM_DESCONTO: total_extras_descontavel += price
                     else: total_extras_nao_descontavel += price

        # Retorna: Custo total mensal só dos extras, parte descontável, parte não descontável
        return total_extras_cost, total_extras_descontavel, total_extras_nao_descontavel

    def atualizar_valores(self, *args):
        if not hasattr(self, 'current_plan') or not self.current_plan or self.current_plan not in PLAN_INFO:
            print("Aviso: Plano atual não definido ou inválido durante atualização de valores.")
            return # Sai se o plano não estiver configurado ainda

        info = PLAN_INFO[self.current_plan]
        is_bling_plan = self.current_plan.startswith("Bling -")
        is_autoatendimento_plano = self.current_plan == "Autoatendimento"
        is_em_branco_plano = self.current_plan == "Em Branco"

        # --- Obter Custos Base e Extras ---
        # base_mensal aqui é o CUSTO MENSAL EFETIVO NO PLANO ANUAL (já com desconto implícito)
        base_mensal_efetivo_anual = info.get("base_mensal", 0.0)
        total_extras_cost, total_extras_descontavel, total_extras_nao_descontavel = self._calcular_extras()

        # --- Calcular Valores Principais ---

        # 1. Base Mensal Sem Fidelidade (Base Efetiva Anual / 0.90)
        #    Isso reverte o desconto padrão de 10% para obter o valor "cheio" da base do plano.
        base_mensal_sem_fidelidade = (base_mensal_efetivo_anual / 0.90) if base_mensal_efetivo_anual > 0.01 else 0.0
        # Caso especial Bling: O valor "base_mensal" já é o anual efetivo, mas podemos ter um 'base_mensal_original' se definido.
        # if is_bling_plan and "base_mensal_original" in info:
        #     base_mensal_sem_fidelidade = info["base_mensal_original"] # Usaria um valor "de" diferente

        # 2. Total Mensal Sem Fidelidade (Base Sem Fid + Custo Total dos Extras)
        total_mensal_sem_fidelidade = base_mensal_sem_fidelidade + total_extras_cost

        # 3. Total Mensal Efetivo Anual - BASE PARA CÁLCULO (Base Efetiva Anual + Custo Total dos Extras)
        #    Este é o valor que seria pago mensalmente no plano anual ANTES de aplicar overrides do usuário.
        total_mensal_efetivo_anual_base_calc = base_mensal_efetivo_anual + total_extras_cost

        # --- Aplicar Overrides ---
        final_mensal_efetivo_anual = 0.0 # Valor final que será exibido como mensal no anual
        final_anual_total = 0.0       # Valor final que será exibido como anual pago adiantado
        desconto_aplicado_percent = 10.0 # Desconto padrão implícito de 10% (base vs base/0.9)

        try:
            if self.user_override_anual_active.get():
                # Usuário digitou o TOTAL ANUAL PAGO ADIANTADO
                edited_total_anual = float(self.valor_anual_editavel.get().replace(",", ".")) # Trata vírgula
                if edited_total_anual < 0: edited_total_anual = 0 # Não permite negativo

                final_anual_total = edited_total_anual
                final_mensal_efetivo_anual = final_anual_total / 12.0

                # Recalcular desconto implícito para exibição
                if total_mensal_sem_fidelidade > 0.01:
                    # Desconto = ( (Valor cheio sem fid) - (Valor efetivo pago mensal no anual) ) / (Valor cheio sem fid)
                    desconto_calc = ((total_mensal_sem_fidelidade - final_mensal_efetivo_anual) / total_mensal_sem_fidelidade) * 100
                    desconto_aplicado_percent = max(0, desconto_calc) # Garante não negativo
                else:
                    desconto_aplicado_percent = 0.0

                # Atualiza o campo de desconto para refletir o valor digitado
                self.desconto_personalizado.set(str(round(desconto_aplicado_percent)))

            elif self.user_override_discount_active.get():
                # Usuário digitou o PERCENTUAL DE DESCONTO
                desc_custom = float(self.desconto_personalizado.get().replace(",", ".")) # Trata vírgula
                if desc_custom < 0: desc_custom = 0 # Não permite negativo
                # if desc_custom > 100: desc_custom = 100 # Pode limitar a 100%?

                desc_dec = desc_custom / 100.0
                desconto_aplicado_percent = desc_custom

                # O desconto customizado incide sobre a parte descontável do valor cheio
                base_sem_fid_mais_extras_descont = base_mensal_sem_fidelidade + total_extras_descontavel
                # O valor final mensal será a parte descontável com o desconto aplicado + a parte não descontável intacta
                final_mensal_efetivo_anual = (base_sem_fid_mais_extras_descont * (1 - desc_dec)) + total_extras_nao_descontavel
                final_anual_total = final_mensal_efetivo_anual * 12.0

                # Atualiza o campo anual para refletir o desconto aplicado
                self.valor_anual_editavel.set(f"{final_anual_total:.2f}")

            else:
                # Cálculo Padrão (sem overrides ativos) - usa o desconto implícito de 10%
                final_mensal_efetivo_anual = total_mensal_efetivo_anual_base_calc # Já inclui extras
                final_anual_total = final_mensal_efetivo_anual * 12.0

                # Calcula o desconto REAL para exibição (pode variar de 10% por causa dos extras não descontáveis)
                if total_mensal_sem_fidelidade > 0.01:
                     desconto_calc = ((total_mensal_sem_fidelidade - final_mensal_efetivo_anual) / total_mensal_sem_fidelidade) * 100
                     desconto_aplicado_percent = max(0, desconto_calc)
                else: desconto_aplicado_percent = 0.0

                # Atualiza campos editáveis com valores padrão calculados
                self.valor_anual_editavel.set(f"{final_anual_total:.2f}")
                # Exibe o desconto calculado (pode não ser exatamente 10)
                self.desconto_personalizado.set(str(round(desconto_aplicado_percent)))

        except ValueError:
             # Se erro na conversão (valor inválido digitado), reverte para o cálculo padrão
             print("Erro de valor na edição manual, revertendo para cálculo padrão.")
             final_mensal_efetivo_anual = total_mensal_efetivo_anual_base_calc
             final_anual_total = final_mensal_efetivo_anual * 12.0
             if total_mensal_sem_fidelidade > 0.01:
                 desconto_calc = ((total_mensal_sem_fidelidade - final_mensal_efetivo_anual) / total_mensal_sem_fidelidade) * 100
                 desconto_aplicado_percent = max(0, desconto_calc)
             else: desconto_aplicado_percent = 0.0
             self.valor_anual_editavel.set(f"{final_anual_total:.2f}")
             self.desconto_personalizado.set(str(round(desconto_aplicado_percent)))
             # Desativa flags de override
             self.user_override_anual_active.set(False)
             self.user_override_discount_active.set(False)


        # --- Calcular Custo Adicional (Treinamento/Implementação) ---
        custo_adicional = 0.0
        label_custo = "Treinamento"
        if is_bling_plan: label_custo = "Implementação"
        elif is_autoatendimento_plano: label_custo = "Setup" # Ou pode ser zero

        # Regra do custo adicional: Aplicável apenas se NÃO for Bling/AA/Em Branco
        if not is_bling_plan and not is_autoatendimento_plano and not is_em_branco_plano:
            limite_custo = 549.90
            # Compara o TOTAL MENSAL SEM FIDELIDADE com o limite
            if total_mensal_sem_fidelidade > 0.01 and total_mensal_sem_fidelidade < limite_custo:
                custo_adicional = limite_custo - total_mensal_sem_fidelidade

        # --- Atualizar Labels da UI ---
        mensal_sem_fid_str = f"{total_mensal_sem_fidelidade:.2f}".replace(".", ",")
        mensal_no_anual_str = f"{final_mensal_efetivo_anual:.2f}".replace(".", ",")
        anual_total_str = f"{final_anual_total:.2f}".replace(".", ",")
        custo_adic_str = f"{custo_adicional:.2f}".replace(".", ",")
        desconto_final_percent = round(desconto_aplicado_percent) # Arredonda para exibição

        self.lbl_plano_mensal_sem_fid.config(text=f"Mensal (Sem Fidelidade): R$ {mensal_sem_fid_str}")
        if custo_adicional > 0.01:
            self.lbl_treinamento.config(text=f"+ Custo {label_custo}: R$ {custo_adic_str}")
            self.lbl_treinamento.pack(pady=(0, 5), anchor="w", padx=15) # Garante que está visível
        else:
            self.lbl_treinamento.pack_forget() # Oculta se for zero

        self.lbl_plano_mensal_no_anual.config(text=f"Mensal (no Plano Anual): R$ {mensal_no_anual_str}")
        self.lbl_plano_anual_total.config(text=f"Anual (Pagamento Único): R$ {anual_total_str}")
        self.lbl_desconto.config(text=f"Desconto Anual Aplicado: {desconto_final_percent}%")

        # --- Armazenar Valores Computados Finais ---
        self.computed_mensal_sem_fidelidade = total_mensal_sem_fidelidade
        self.computed_mensal_efetivo_anual = final_mensal_efetivo_anual
        self.computed_anual_total = final_anual_total
        self.computed_desconto_percent = desconto_final_percent # Armazena arredondado
        self.computed_custo_adicional = custo_adicional

    def montar_lista_modulos(self):
        """ Monta a string formatada com a lista de módulos ativos para a proposta. """
        linhas = []
        info = PLAN_INFO.get(self.current_plan, {})
        mandatory = info.get("mandatory", [])
        modulos_ja_listados = set() # Para evitar duplicatas

        # 1. Itens de Spinbox (se > 0 ou > mínimo)
        pdv_val = self.spin_pdv_var.get()
        if pdv_val > 0:
            linhas.append(f"{pdv_val}x PDV - Frente de Caixa")
            modulos_ja_listados.add("PDV - Frente de Caixa")

        usr_val = self.spin_users_var.get()
        if usr_val > 0:
            linhas.append(f"{usr_val}x Usuários")
            modulos_ja_listados.add("Usuários")

        smart_tef_val = self.spin_smart_tef_var.get()
        if smart_tef_val > 0:
            # Adiciona nota sobre limite se for plano Gestão
            extra_info = " (Limite: 3)" if self.current_plan == "Gestão" else ""
            # Se for Performance, indica que já estão incluídos
            if self.current_plan == "Performance" and smart_tef_val == 3:
                 extra_info = " (Inclusos no plano)"

            linhas.append(f"{smart_tef_val}x Smart TEF{extra_info}")
            modulos_ja_listados.add("Smart TEF")

        taa_val = self.spin_terminais_aa_var.get()
        if taa_val > 0:
             # Se for plano Autoatendimento e for 1, indica que está incluso
             extra_info = ""
             if self.current_plan == "Autoatendimento" and taa_val == 1:
                 extra_info = " (Incluso no plano)"
             linhas.append(f"{taa_val}x Terminais Autoatendimento{extra_info}")
             modulos_ja_listados.add("Terminais Autoatendimento")


        # 2. Módulos Mandatórios (que ainda não foram listados pelos spinboxes)
        for m in mandatory:
            if m not in modulos_ja_listados:
                 # Adiciona "1x " para clareza, exceto para itens como Notas Fiscais
                 prefix = "1x "
                 if "Notas Fiscais" in m or "Suporte" in m or "Relatório" in m:
                      prefix = ""
                 linhas.append(f"{prefix}{m}")
                 modulos_ja_listados.add(m)

        # 3. Módulos Opcionais (Checkboxes marcados que não são mandatórios)
        for m, var_m in self.modules.items():
            # Considera apenas checkboxes visíveis e marcados
            if m in self.check_buttons and var_m.get() == 1:
                if m not in mandatory and m not in modulos_ja_listados:
                     linhas.append(f"1x {m}")
                     modulos_ja_listados.add(m) # Evita duplicar se estiver em allowed e for marcado

        # 4. Formatação Final
        # Remove duplicatas (embora a lógica acima deva prevenir) e ordena se desejado
        unique_mods = []
        for item in linhas:
            if item not in unique_mods:
                unique_mods.append(item)
        # unique_mods.sort() # Descomente se quiser ordenar alfabeticamente

        montagem = "\n".join(f"•    {m}" for m in unique_mods)
        return montagem

    def gerar_dados_proposta(self, nome_closer, cel_closer, email_closer):
        """ Gera o dicionário de dados para preencher o slide da proposta. """
        # Nome do Plano: Usa o editado se preenchido, senão o selecionado
        nome_plano_selecionado = self.current_plan
        nome_plano_editado = self.nome_plano_var.get().strip()
        nome_plano_final = nome_plano_editado if nome_plano_editado else nome_plano_selecionado

        # --- Valores Formatados ---
        # Usar os valores já calculados e armazenados em atualizar_valores()
        mensal_sem_fid_val = self.computed_mensal_sem_fidelidade
        mensal_efetivo_anual_val = self.computed_mensal_efetivo_anual
        anual_total_val = self.computed_anual_total
        custo_adicional_val = self.computed_custo_adicional
        desconto_percent_val = self.computed_desconto_percent # Já arredondado

        # Strings formatadas com R$ e vírgula decimal
        mensal_sem_fid_str = f"R$ {mensal_sem_fid_val:.2f}".replace(".", ",")
        mensal_efetivo_anual_str = f"R$ {mensal_efetivo_anual_val:.2f}".replace(".", ",")
        anual_total_str = f"R$ {anual_total_val:.2f}".replace(".", ",")
        custo_adicional_str = f"R$ {custo_adicional_val:.2f}".replace(".", ",")

        # Adicionar custo adicional ao mensal sem fidelidade na string (se aplicável)
        label_custo = "Treinamento"
        if self.current_plan.startswith("Bling"): label_custo = "Implementação"
        elif self.current_plan == "Autoatendimento": label_custo = "Setup" # Ou pode ser "Incluso"

        str_mensal_completa = mensal_sem_fid_str
        if custo_adicional_val > 0.01:
            str_mensal_completa += f" + {custo_adicional_str} ({label_custo})"

        # --- Definir Suporte ---
        # Regra simplificada: Estendido se Performance ou opcional marcado, senão Regular/Chat
        tipo_suporte = "Regular" # Default
        horario_suporte = "09:00 às 17:00 de Segunda a Sexta-feira"

        suporte_chat_ativo = self.modules.get("Suporte Técnico - Via chat", tk.IntVar()).get() == 1
        suporte_estendido_ativo = self.modules.get("Suporte Técnico - Estendido", tk.IntVar()).get() == 1

        if self.current_plan == "Performance" or suporte_estendido_ativo:
             tipo_suporte = "Estendido"
             horario_suporte = "09:00 às 22:00 Seg-Sex & 11:00 às 21:00 Sab-Dom"
        elif suporte_chat_ativo: # Se não for estendido mas tiver chat (Plano Gestão)
             tipo_suporte = "Chat Incluso"
             horario_suporte = "09:00 às 22:00 Seg-Sex & 11:00 às 21:00 Sab-Dom"
        # PDV Básico fica com o default "Regular"

        # --- Montar Lista de Módulos ---
        montagem_formatada = self.montar_lista_modulos()

        # --- Calcular Economia Anual ---
        economia_str = ""
        # Custo total se pagasse mensalmente por 12 meses = (Mensal Sem Fidelidade * 12) + Custo Adicional (pago uma vez)
        custo_total_mensalizado = (mensal_sem_fid_val * 12) + custo_adicional_val
        # Custo total pagando anualmente adiantado
        custo_total_anualizado = anual_total_val # Já é o total pago adiantado

        # Compara os dois totais anuais
        if custo_total_mensalizado > custo_total_anualizado + 0.01: # Adiciona margem para evitar R$ 0,00 por arredondamento
             economia_val = custo_total_mensalizado - custo_total_anualizado
             econ_str = f"{economia_val:.2f}".replace(".", ",")
             economia_str = f"Economia de R$ {econ_str} no plano anual" # Texto ajustado

        # --- Montar Dicionário Final para Placeholders ---
        # **IMPORTANTE:** As chaves aqui DEVEM CORRESPONDER EXATAMENTE aos placeholders no seu PPTX.
        # Ex: se no PPTX está {NomeCliente}, a chave deve ser "NomeCliente".
        # Ajuste as chaves abaixo conforme seu template.
        dados = {
            "montagem_do_plano": montagem_formatada,
            "plano_mensal_sem_fidelidade": str_mensal_completa,      # Ex: R$ 110,00 + R$ 439,90 (Treinamento) ou R$ 600,00
            "plano_mensal_no_anual": mensal_efetivo_anual_str,    # Ex: R$ 99,00
            "plano_anual_total": anual_total_str,            # Ex: R$ 1188,00
            "custo_treinamento": custo_adicional_str if custo_adicional_val > 0.01 else "Incluso", # Ex: R$ 439,90 ou Incluso
            "desconto_aplicado": f"{desconto_percent_val}%",     # Ex: 10%
            "nome_do_plano": nome_plano_final,               # Ex: Performance ou Plano Super Cliente
            "tipo_de_suporte": tipo_suporte,                 # Ex: Estendido
            "horario_de_suporte": horario_suporte,           # Ex: 09:00 às 22:00...
            "validade_proposta": self.validade_proposta_var.get(), # Ex: 05/05/2025
            "nome_closer": nome_closer,                      # Ex: João Silva
            "celular_closer": cel_closer,                  # Ex: (41) 99999-8888
            "email_closer": email_closer,                    # Ex: joao.silva@email.com
            "nome_cliente": self.nome_cliente_var.get(),     # Ex: Padaria Pão Quente
            "economia_anual": economia_str                   # Ex: Economia de R$ 120,00 no plano anual
        }
        # Adicionar placeholders extras se necessário (ex: data atual)
        dados["data_geracao"] = date.today().strftime("%d/%m/%Y")

        return dados


# --- Funções de Geração de PPTX (Proposta e Material) ---

# *** MAPEAMENTO DE MÓDULOS PARA SLIDES ***
# Ajuste as chaves (placeholders no PPTX) e valores (nomes dos módulos no Python)
# Conforme os seus arquivos .pptx de template.
# Use 'None' para slides que devem sempre aparecer (se tiverem um placeholder como {slide_sempre}).
# Use um set de módulos se múltiplos módulos ativam o mesmo slide.
MAPEAMENTO_MODULOS_MATERIAL = {
    "slide_sempre": None, # Exemplo de slide que sempre fica
    "slide_bling": {"Bling - Básico", "Bling - Com Estoque em Grade"}, # Slide específico do Bling
    "check_sistema_kds": "Relatório KDS",
    "check_Hub_de_Delivery": "Hub de Delivery",
    "check_integracao_api": "Integração API",
    "check_integracao_tap": "Integração Tap",
    "check_controle_de_mesas": "Controle de Mesas",
    "check_Delivery": "Delivery",
    "check_producao": "Produção",
    "check_Estoque_em_Grade": "Estoque em Grade",
    "check_Facilita_NFE": "Facilita NFE",
    "check_Importacao_de_xml": "Importação de XML",
    "check_conciliacao_bancaria": "Conciliação Bancária",
    "check_contratos_de_cartoes": "Contratos de cartões e outros",
    "check_ordem_de_servico": "Ordem de Serviço",
    "check_relatorio_dinamico": "Relatório Dinâmico",
    "check_programa_de_fidelidade": "Programa de Fidelidade",
    "check_business_intelligence": "Business Intelligence (BI)",
    "check_smartmenu": "Smart Menu",
    "check_backup_real_time": "Backup Realtime",
    "check_att_tempo_real": "Atualização em tempo real",
    "check_promocao": "Promoções",
    "check_marketing": "Marketing",
    "placeholder_pdv": "PDV - Frente de Caixa", # Associado a ter PDV > 0
    "placeholder_smarttef": "Smart TEF", # Associado a ter SmartTEF > 0
    "placeholder_tef": "TEF", # Associado a ter TEF (opcional) > 0
    "placeholder_autoatendimento": "Terminais Autoatendimento", # Associado a ter TAA > 0
    "placeholder_cardapio_digital": "Cardápio digital",
    "placeholder_app_gestao_cplug": "App Gestão CPlug",
    "check_delivery_direto_vip": "Delivery Direto VIP",
    "check_delivery_direto_profissional": "Delivery Direto Profissional",
    "placeholder_painel_senha_tv": "Painel Senha TV",
    "placeholder_painel_senha_mobile": "Painel Senha Mobile",
    "placeholder_suporte_chat": "Suporte Técnico - Via chat",
    "placeholder_suporte_estendido": "Suporte Técnico - Estendido",
    "placeholder_notas_fiscais": {"Notas Fiscais Ilimitadas", "30 Notas Fiscais"}, # Qualquer tipo de NF ativa
}

# Mapeamento para a PROPOSTA (se for diferente do material)
# Se a proposta não tiver slides condicionais, pode deixar vazio ou usar apenas 'slide_sempre'
MAPEAMENTO_MODULOS_PROPOSTA = {
    "slide_sempre": None,
    # Adicione aqui placeholders específicos da PROPOSTA que condicionam slides
    # Ex: "slide_condicional_proposta": "ModuloX",
    # Se a proposta sempre mostra todos os slides e só muda o conteúdo,
    # este dicionário pode ficar apenas com 'slide_sempre': None.
}


def _processar_geracao_pptx(tipo_arquivo, pptx_template_path, lista_abas, nome_closer, cel_closer, email_closer, mapeamento_slides):
    """Função interna para gerar Proposta ou Material Técnico."""

    if not os.path.exists(pptx_template_path):
        showerror("Erro", f"Arquivo template '{os.path.basename(pptx_template_path)}' não encontrado!")
        return None

    try:
        prs = Presentation(pptx_template_path)
    except Exception as e:
        showerror("Erro", f"Falha ao abrir '{os.path.basename(pptx_template_path)}': {e}")
        return None

    if not lista_abas:
        showerror("Erro", f"Não há abas ativas para gerar {tipo_arquivo}.")
        return None

    # --- 1. Coletar dados e módulos ativos ---
    # Usa os dados da PRIMEIRA aba para placeholders globais e remoção de slides
    primeira_aba = lista_abas[0]
    dados_globais = primeira_aba.gerar_dados_proposta(nome_closer, cel_closer, email_closer)

    # Coleta todos os módulos ativos de TODAS as abas para decidir quais slides manter
    modulos_ativos_geral = set()
    planos_usados_geral = set()
    for aba in lista_abas:
        planos_usados_geral.add(aba.current_plan)
        info_aba = PLAN_INFO.get(aba.current_plan, {})
        mandatory_aba = info_aba.get("mandatory", [])
        # Adiciona mandatórios
        for mod in mandatory_aba: modulos_ativos_geral.add(mod)
        # Adiciona opcionais marcados
        for nome_mod, var_mod in aba.modules.items():
            if var_mod.get() == 1: # Não precisa checar se é mandatório aqui, o set cuida disso
                 modulos_ativos_geral.add(nome_mod)
        # Adiciona itens de spinbox (se > 0)
        if aba.spin_pdv_var.get() > 0: modulos_ativos_geral.add("PDV - Frente de Caixa")
        if aba.spin_users_var.get() > 0: modulos_ativos_geral.add("Usuários")
        if aba.spin_smart_tef_var.get() > 0: modulos_ativos_geral.add("Smart TEF")
        if aba.spin_terminais_aa_var.get() > 0: modulos_ativos_geral.add("Terminais Autoatendimento")
        if aba.modules.get("TEF", tk.IntVar()).get() == 1: modulos_ativos_geral.add("TEF") # Checkbox TEF
        # Adicionar outros módulos controlados por spinbox ou UI se relevante para os slides

    print(f"--- {tipo_arquivo} ---")
    print(f"Planos usados: {planos_usados_geral}")
    print(f"Módulos ativos (geral): {modulos_ativos_geral}")


    # --- 2. Decidir quais slides manter ---
    keep_slides = set()
    num_slides_original = len(prs.slides)
    print(f"Analisando {num_slides_original} slides do template...")

    for i, slide in enumerate(prs.slides):
        slide_mantido = False
        textos_no_slide = [] # Para debug

        # Verifica o texto em todas as shapes do slide
        for shape in slide.shapes:
             if slide_mantido: break # Otimização: se já decidiu manter, vai pro próximo slide
             if shape.has_text_frame:
                 for paragraph in shape.text_frame.paragraphs:
                     if slide_mantido: break
                     for run in paragraph.runs:
                         if slide_mantido: break
                         txt_run = run.text.strip()
                         if not txt_run: continue
                         textos_no_slide.append(txt_run) # Guarda texto para debug

                         # Verifica se algum placeholder do mapeamento está no texto
                         for placeholder, modulo_mapeado in mapeamento_slides.items():
                             if placeholder in txt_run:
                                 if modulo_mapeado is None: # Placeholder {slide_sempre} por exemplo
                                     slide_mantido = True; break
                                 elif isinstance(modulo_mapeado, set): # Se o valor for um conjunto de módulos
                                     # Mantém se QUALQUER um dos módulos do conjunto estiver ativo
                                     if any(m in modulos_ativos_geral for m in modulo_mapeado):
                                         slide_mantido = True; break
                                     # Caso especial: Se for placeholder de plano Bling, checa planos_usados_geral
                                     if placeholder == "slide_bling" and any(p.startswith("Bling -") for p in planos_usados_geral):
                                          slide_mantido = True; break
                                 elif isinstance(modulo_mapeado, str): # Se o valor for uma string (um módulo)
                                     if modulo_mapeado in modulos_ativos_geral:
                                         slide_mantido = True; break
                                 # Se chegou aqui, placeholder foi encontrado mas condição não bateu

        # Adiciona ao conjunto se deve ser mantido
        if slide_mantido:
            keep_slides.add(i)
        # else: # Debug Opcional: mostrar texto de slides removidos
        #     print(f"  -> Removendo Slide {i+1}. Textos encontrados: {textos_no_slide[:5]}...")


    # --- 3. Remover slides não mantidos ---
    num_slides_remover = num_slides_original - len(keep_slides)
    print(f"Decisão: Manter {len(keep_slides)} slides, remover {num_slides_remover}.")

    # Itera de trás para frente para não afetar índices dos slides restantes
    if num_slides_remover > 0:
        for idx in reversed(range(num_slides_original)):
            if idx not in keep_slides:
                try:
                    # Método robusto para remover slide
                    rId = prs.slides._sldIdLst[idx].rId
                    prs.part.drop_rel(rId) # Remove relacionamento
                    del prs.slides._sldIdLst[idx] # Remove da lista de slides
                    # print(f"   Slide {idx+1} removido.")
                except Exception as e:
                    # Pode dar erro se o slide já foi removido ou tiver alguma inconsistência
                    print(f"Aviso: Não foi possível remover slide {idx+1}. Erro: {e}")
        print(f"Remoção concluída. Slides restantes: {len(prs.slides)}")
    else:
        print("Nenhum slide para remover.")


    # --- 4. Substituir placeholders nos slides restantes ---
    print("Substituindo placeholders...")
    for slide in prs.slides:
        # Usa a função corrigida com os dados globais (da primeira aba)
        substituir_placeholders_no_slide(slide, dados_globais)


    # --- 5. Salvar o arquivo final ---
    nome_cliente_safe = dados_globais.get("nome_cliente", "SemNome").replace("/", "-").replace("\\", "-")
    hoje_str = date.today().strftime("%d-%m-%Y")
    prefixo_arquivo = "Proposta ConnectPlug" if tipo_arquivo == "Proposta" else "Material Tecnico ConnectPlug"
    nome_arquivo_final = f"{prefixo_arquivo} - {nome_cliente_safe} - {hoje_str}.pptx"

    try:
        prs.save(nome_arquivo_final)
        showinfo("Sucesso", f"{tipo_arquivo} gerada com sucesso:\n{nome_arquivo_final}")
        return nome_arquivo_final
    except PermissionError as e:
         showerror("Erro de Permissão", f"Não foi possível salvar '{nome_arquivo_final}'.\nVerifique se o arquivo não está aberto ou se você tem permissão para escrever na pasta.\nErro: {e}")
         return None
    except Exception as e:
        showerror("Erro ao Salvar", f"Falha desconhecida ao salvar '{nome_arquivo_final}':\n{e}")
        return None


# Funções wrapper que chamam a função interna
def gerar_proposta(lista_abas, nome_closer, celular_closer, email_closer):
    return _processar_geracao_pptx(
        tipo_arquivo="Proposta",
        pptx_template_path="Proposta Comercial ConnectPlug.pptx",
        lista_abas=lista_abas,
        nome_closer=nome_closer,
        cel_closer=celular_closer,
        email_closer=email_closer,
        mapeamento_slides=MAPEAMENTO_MODULOS_PROPOSTA # Usa mapeamento da proposta
    )

def gerar_material(lista_abas, nome_closer, celular_closer, email_closer):
     return _processar_geracao_pptx(
        tipo_arquivo="Material Técnico",
        pptx_template_path="Material Tecnico ConnectPlug.pptx",
        lista_abas=lista_abas,
        nome_closer=nome_closer,
        cel_closer=celular_closer,
        email_closer=email_closer,
        mapeamento_slides=MAPEAMENTO_MODULOS_MATERIAL # Usa mapeamento do material
    )


# --- Google Drive / Auth / Upload (SEM ALTERAÇÕES NECESSÁRIAS NA LÓGICA PRINCIPAL) ---
SCOPES = ['https://www.googleapis.com/auth/drive']
CLIENT_SECRET_URL = "https://github.com/DevRGS/Gerador/raw/refs/heads/main/config/client_secret_788265418970-ur6f189oqvsttseeg6g77fegt0su67dj.apps.googleusercontent.com.json"
CLIENT_SECRET_LOCAL_FILE = "client_secret_temp.json"
TOKEN_FILE = 'token.json'

def baixar_client_secret_remoto():
    """Baixa o client_secret.json do GitHub se não existir localmente."""
    if not os.path.exists(CLIENT_SECRET_LOCAL_FILE):
        print(f"Baixando {CLIENT_SECRET_LOCAL_FILE} do GitHub...")
        try:
            r = requests.get(CLIENT_SECRET_URL, timeout=15)
            r.raise_for_status()
            with open(CLIENT_SECRET_LOCAL_FILE, "w", encoding="utf-8") as f:
                 f.write(r.text)
            print("Client secret baixado com sucesso.")
        except requests.exceptions.Timeout:
             showerror("Erro de Rede", f"Tempo esgotado ao tentar baixar {CLIENT_SECRET_LOCAL_FILE}.")
             raise Exception("Timeout ao baixar client_secret.")
        except requests.exceptions.RequestException as e:
             showerror("Erro de Rede", f"Não foi possível baixar o {CLIENT_SECRET_LOCAL_FILE}:\n{e}")
             raise Exception(f"Erro ao baixar o client_secret.json: {e}")
    return CLIENT_SECRET_LOCAL_FILE

def get_gdrive_service():
    """Autentica (ou usa token salvo) e retorna o serviço do Google Drive."""
    creds = None
    client_secret_file = None
    try:
        client_secret_file = baixar_client_secret_remoto()
    except Exception as e:
        showerror("Erro Crítico", f"Falha ao obter client_secret.\nVerifique a conexão com a internet.\n{e}")
        return None # Não pode continuar sem client secret

    # Tenta carregar token existente
    if os.path.exists(TOKEN_FILE):
        try:
            with open(TOKEN_FILE, 'rb') as token:
                creds = pickle.load(token)
            print("Token local carregado.")
        except (pickle.UnpicklingError, EOFError, FileNotFoundError, Exception) as e:
             print(f"Erro ao carregar token ({e}). Removendo arquivo token corrompido/inválido.")
             try:
                 os.remove(TOKEN_FILE)
             except OSError:
                 pass
             creds = None

    # Valida credenciais ou re-autentica
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("Token expirado, tentando renovar...")
            try:
                creds.refresh(Request())
                print("Token renovado com sucesso.")
                # Salva o token renovado
                try:
                    with open(TOKEN_FILE, 'wb') as token: pickle.dump(creds, token)
                except Exception as e: print(f"Aviso: Não foi possível salvar token renovado: {e}")
            except Exception as e:
                print(f"Erro ao renovar token: {e}. Reautenticando...")
                if os.path.exists(TOKEN_FILE):
                    try: os.remove(TOKEN_FILE); print("Token antigo removido.")
                    except OSError: pass
                creds = None # Força re-autenticação
        else:
             # Se não tem credencial ou não pôde renovar, faz fluxo de autenticação
             if os.path.exists(TOKEN_FILE): # Remove token inválido se chegou aqui sem creds
                 try: os.remove(TOKEN_FILE); print("Removendo token inválido existente.")
                 except OSError: pass
             creds = None

        if not creds: # Se ainda não tem credencial válida, roda o fluxo
             print("Nenhuma credencial válida encontrada. Iniciando fluxo de autenticação...")
             try:
                flow = InstalledAppFlow.from_client_secrets_file(client_secret_file, SCOPES)
                # Tenta abrir navegador, com fallback para console se GUI não disponível
                creds = flow.run_local_server(port=0, open_browser=True)
                print("Autenticação bem-sucedida.")
                # Salva as novas credenciais
                try:
                    with open(TOKEN_FILE, 'wb') as token: pickle.dump(creds, token)
                    print(f"Novas credenciais salvas em {TOKEN_FILE}.")
                except Exception as e: print(f"Aviso: Não foi possível salvar o novo token.json: {e}")
             except FileNotFoundError: # Caso o client_secret_temp não exista mais
                 showerror("Erro de Autenticação", "Arquivo client_secret não encontrado durante a autenticação.")
                 return None
             except Exception as e:
                 showerror("Erro de Autenticação", f"Falha durante o fluxo de autenticação:\n{e}")
                 return None

    # Constrói e retorna o serviço
    try:
        service = build('drive', 'v3', credentials=creds)
        print("Serviço Google Drive construído com sucesso.")
        return service
    except Exception as e:
        showerror("Erro Google API", f"Falha ao construir serviço Google Drive:\n{e}")
        return None

def upload_pptx_and_export_to_pdf(local_pptx_path):
    """Faz upload do PPTX para o Google Drive (como Google Slides) e exporta como PDF."""
    if not os.path.exists(local_pptx_path):
        showerror("Erro", f"Arquivo '{os.path.basename(local_pptx_path)}' não encontrado para upload.")
        return None # Retorna None para indicar falha

    service = get_gdrive_service()
    if not service:
        showerror("Erro Google Drive", "Não foi possível conectar ao Google Drive para gerar o PDF.")
        return None # Retorna None

    pdf_output_name = local_pptx_path.replace(".pptx", ".pdf")
    base_name = os.path.basename(local_pptx_path)
    file_id = None # Para garantir que file_id exista no bloco finally

    try:
        print(f"Iniciando upload de '{base_name}' para o Google Drive...")
        # Metadata: Converte para Google Slides no upload
        file_metadata = {'name': base_name, 'mimeType': 'application/vnd.google-apps.presentation'}
        # Media: Conteúdo do arquivo PPTX
        media = MediaFileUpload(local_pptx_path,
                                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                                resumable=True)

        # Executa o upload
        uploaded_file = service.files().create(body=file_metadata,
                                             media_body=media,
                                             fields='id').execute()
        file_id = uploaded_file.get('id')
        if not file_id:
            raise Exception("Falha no upload (ID do arquivo não retornado).")

        print(f"Upload concluído. ID do arquivo no Drive: {file_id}. Iniciando exportação para PDF...")

        # Exporta o arquivo (agora Google Slides) para PDF
        request = service.files().export_media(fileId=file_id, mimeType='application/pdf')
        fh = io.BytesIO() # Buffer em memória para baixar o PDF
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            if status:
                 print(f"Download PDF: {int(status.progress() * 100)}%")

        # Salva o PDF baixado localmente
        with open(pdf_output_name, 'wb') as f:
            f.write(fh.getvalue())

        print(f"PDF gerado com sucesso: '{pdf_output_name}'.")
        showinfo("Google Drive", f"PDF gerado com sucesso:\n'{os.path.basename(pdf_output_name)}'.")
        return pdf_output_name # Retorna o nome do PDF gerado

    except Exception as e:
        showerror("Erro Google Drive", f"Erro durante o processo de upload/conversão para PDF:\n{e}")
        return None # Retorna None em caso de erro
    finally:
        # Tenta deletar o arquivo temporário do Drive (que foi convertido para GSlides)
        if file_id:
            try:
                print(f"Tentando deletar arquivo temporário do Drive (ID: {file_id})...")
                service.files().delete(fileId=file_id).execute()
                print("Arquivo temporário deletado do Google Drive.")
            except Exception as delete_err:
                # Não é crítico se falhar, apenas avisa
                print(f"Aviso: Falha ao deletar arquivo temporário do Google Drive: {delete_err}")


# --- MainApp (Interface Gráfica Principal) ---
class MainApp(ttkb.Window):
    def __init__(self):
        # Define o tema ANTES de chamar super().__init__
        # try:
        #     # Tenta usar um tema visualmente agradável do ttkbootstrap
        #     style = ttkb.Style(theme='litera') # 'litera', 'flatly', 'journal', 'pulse', etc.
        #     super().__init__(themename=style.theme.name)
        # except Exception:
        #     # Fallback para tema padrão se falhar
        #     super().__init__()

        super().__init__(themename="litera") # Escolha um tema: litera, cosmo, flatly, journal, etc.

        self.title("Gerador de Propostas ConnectPlug v2.2 - Corrigido") # Versão
        # Define um tamanho inicial razoável
        self.geometry("1200x850")
        # Permite redimensionamento
        self.resizable(True, True)

        # Variáveis do Vendedor
        self.nome_closer_var = tk.StringVar()
        self.celular_closer_var = tk.StringVar()
        self.email_closer_var = tk.StringVar()

        # Variáveis compartilhadas entre as Abas
        self.nome_cliente_var_shared = tk.StringVar(value="") # Começa vazio
        self.validade_proposta_var_shared = tk.StringVar(value=date.today().strftime("%d/%m/%Y")) # Default hoje

        # Carrega config do vendedor e define ação ao fechar
        carregar_config(self.nome_closer_var, self.celular_closer_var, self.email_closer_var)
        self.protocol("WM_DELETE_WINDOW", self.on_close)

        # --- Layout Principal ---
        # Top Bar (Dados Vendedor e Botão Nova Aba)
        top_bar = ttkb.Frame(self, padding=5)
        top_bar.pack(side="top", fill="x")

        ttkb.Label(top_bar, text="Vendedor:").pack(side="left", padx=(0, 2))
        ttkb.Entry(top_bar, textvariable=self.nome_closer_var, width=20).pack(side="left", padx=(0, 10))
        ttkb.Label(top_bar, text="Celular:").pack(side="left", padx=(0, 2))
        ttkb.Entry(top_bar, textvariable=self.celular_closer_var, width=15).pack(side="left", padx=(0, 10))
        ttkb.Label(top_bar, text="Email:").pack(side="left", padx=(0, 2))
        ttkb.Entry(top_bar, textvariable=self.email_closer_var, width=25).pack(side="left", padx=(0, 10))

        # Botão para adicionar nova aba, alinhado à direita
        self.btn_add = ttkb.Button(top_bar, text="+ Nova Aba", command=self.add_aba, bootstyle="success")
        # self.btn_add.pack(side="right", padx=5) # Descomente se quiser à direita
        self.btn_add.pack(side="left", padx=(20, 5)) # Ou à esquerda depois dos dados


        # Notebook (para as abas dos planos)
        self.notebook = ttkb.Notebook(self, padding=5)
        self.notebook.pack(fill="both", expand=True, pady=5)

        # Bottom Frame (Botões de Geração)
        bot_frame = ttkb.Frame(self, padding=10)
        bot_frame.pack(side="bottom", fill="x")

        ttkb.Button(bot_frame, text="Gerar Proposta + PDF", command=self.on_gerar_proposta, bootstyle="primary").pack(side="left", padx=5)
        ttkb.Button(bot_frame, text="Gerar Material + PDF", command=self.on_gerar_mat_tecnico, bootstyle="info").pack(side="left", padx=5)
        ttkb.Button(bot_frame, text="Gerar TUDO + PDF", command=self.on_gerar_tudo, bootstyle="secondary").pack(side="left", padx=5)


        # --- Inicialização ---
        self.abas_criadas = {} # Dicionário para rastrear abas {indice: widget_frame}
        self.ultimo_indice = 0 # Contador para nomes das abas
        self.add_aba() # Cria a primeira aba ao iniciar

        # Baixa os templates PPTX se não existirem
        # (Colocado após criar a UI principal para não bloquear a janela)
        self.after(100, self.baixar_templates_necessarios) # Atraso pequeno para UI aparecer


    def baixar_templates_necessarios(self):
         """Função chamada para baixar os templates após a UI iniciar."""
         print("Verificando templates PPTX...")
         try:
            baixar_arquivo_if_needed("Proposta Comercial ConnectPlug.pptx", "https://github.com/DevRGS/Gerador/raw/refs/heads/main/assets/Proposta%20Comercial%20ConnectPlug.pptx")
            baixar_arquivo_if_needed("Material Tecnico ConnectPlug.pptx", "https://github.com/DevRGS/Gerador/raw/refs/heads/main/assets/Material%20Tecnico%20ConnectPlug.pptx")
         except Exception as e:
             # showerror já é chamado dentro de baixar_arquivo_if_needed em caso de falha
             print(f"Erro ao baixar templates: {e}")
             # Considerar se a aplicação deve fechar se não conseguir baixar


    def on_close(self):
        """Salva a configuração do vendedor antes de fechar."""
        print("Salvando configuração do vendedor...")
        salvar_config(self.nome_closer_var.get(), self.celular_closer_var.get(), self.email_closer_var.get())
        print("Fechando aplicação.")
        self.destroy()

    def add_aba(self):
        """Adiciona uma nova aba (PlanoFrame) ao Notebook."""
        if len(self.abas_criadas) >= MAX_ABAS:
            showinfo("Limite Atingido", f"Máximo de {MAX_ABAS} abas simultâneas atingido.")
            return

        self.ultimo_indice += 1
        idx = self.ultimo_indice
        print(f"Adicionando Aba Plano {idx}...")

        # Cria a nova aba (PlanoFrame) passando as variáveis compartilhadas
        frame_aba = PlanoFrame(
            self.notebook,
            idx,
            self.nome_cliente_var_shared,
            self.validade_proposta_var_shared,
            self.fechar_aba # Passa a função de callback para fechar
        )

        # Adiciona ao notebook e ao dicionário de rastreamento
        self.notebook.add(frame_aba, text=f"Plano {idx}")
        self.abas_criadas[idx] = frame_aba
        self.notebook.select(frame_aba) # Seleciona a aba recém-criada

        # Atualiza estado do botão de adicionar
        if len(self.abas_criadas) >= MAX_ABAS:
            self.btn_add.config(state="disabled")

    def fechar_aba(self, indice):
        """Remove a aba especificada pelo índice."""
        if indice in self.abas_criadas:
            frame_aba = self.abas_criadas[indice]
            try:
                print(f"Fechando Aba Plano {indice}...")
                self.notebook.forget(frame_aba) # Remove do Notebook visualmente
                del self.abas_criadas[indice] # Remove do dicionário de rastreamento
            except tk.TclError as e:
                 # Pode ocorrer se a aba já foi destruída ou é inválida
                 print(f"Erro ao tentar fechar aba {indice}: {e}")
                 # Garante remoção do dicionário mesmo se forget() falhar
                 if indice in self.abas_criadas: del self.abas_criadas[indice]

            # Reabilita botão de adicionar se abaixo do limite
            if len(self.abas_criadas) < MAX_ABAS:
                self.btn_add.config(state="normal")

            # Se não houver mais abas, cria uma nova automaticamente
            if not self.abas_criadas:
                print("Nenhuma aba restante, adicionando uma nova.")
                self.add_aba()
        else:
             print(f"Aviso: Tentativa de fechar aba com índice inválido: {indice}")


    def get_abas_ativas(self):
        """Retorna uma lista dos widgets PlanoFrame das abas atualmente abertas, ordenadas por índice."""
        indices_ativos = sorted(self.abas_criadas.keys())
        return [self.abas_criadas[idx] for idx in indices_ativos]

    def _validar_dados_basicos(self):
        """Verifica se os dados essenciais do vendedor e cliente estão preenchidos."""
        erros = []
        if not self.nome_closer_var.get(): erros.append("Nome do Vendedor")
        if not self.celular_closer_var.get(): erros.append("Celular do Vendedor")
        if not self.email_closer_var.get(): erros.append("Email do Vendedor")
        if not self.nome_cliente_var_shared.get(): erros.append("Nome do Cliente")
        # Validade não é estritamente obrigatória, mas pode adicionar se quiser

        if erros:
            msg_erro = "Por favor, preencha os seguintes campos antes de gerar:\n- " + "\n- ".join(erros)
            showerror("Dados Incompletos", msg_erro)
            return False
        return True

    def _executar_geracao_e_pdf(self, funcao_geracao_pptx, *args):
         """Função auxiliar para gerar PPTX e depois o PDF."""
         pptx_file = funcao_geracao_pptx(*args)
         pdf_file = None
         if pptx_file and os.path.exists(pptx_file):
             try:
                 # Tenta gerar o PDF a partir do PPTX recém-criado
                 pdf_file = upload_pptx_and_export_to_pdf(pptx_file)
             except Exception as e:
                  print(f"Erro ao gerar PDF para '{os.path.basename(pptx_file)}': {e}")
                  showerror("Erro PDF", f"Falha ao gerar o PDF para '{os.path.basename(pptx_file)}'.\nVerifique a conexão e as permissões do Google Drive.\nErro: {e}")
         return pptx_file, pdf_file # Retorna nomes dos arquivos (ou None se falhar)

    def on_gerar_proposta(self):
        """Gera a Proposta Comercial (PPTX e PDF)."""
        abas_ativas = self.get_abas_ativas()
        if not abas_ativas: showerror("Erro", "Nenhuma aba ativa para gerar a proposta."); return
        if not self._validar_dados_basicos(): return

        print("Iniciando geração da Proposta Comercial...")
        self._executar_geracao_e_pdf(
            gerar_proposta, # Função que gera o PPTX da proposta
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        print("Processo de geração da Proposta concluído.")


    def on_gerar_mat_tecnico(self):
        """Gera o Material Técnico (PPTX e PDF)."""
        abas_ativas = self.get_abas_ativas()
        if not abas_ativas: showerror("Erro", "Nenhuma aba ativa para gerar o material técnico."); return
        if not self._validar_dados_basicos(): return # Valida cliente também

        print("Iniciando geração do Material Técnico...")
        self._executar_geracao_e_pdf(
            gerar_material, # Função que gera o PPTX do material
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        print("Processo de geração do Material Técnico concluído.")


    def on_gerar_tudo(self):
        """Gera AMBOS os documentos (Proposta e Material), cada um com seu PDF."""
        abas_ativas = self.get_abas_ativas()
        if not abas_ativas: showerror("Erro", "Nenhuma aba ativa para gerar os documentos."); return
        if not self._validar_dados_basicos(): return

        print("--- Iniciando Geração de TUDO (Proposta e Material) ---")

        # 1. Gerar Proposta + PDF
        print("\n[1/2] Gerando Proposta...")
        _, pdf_prop = self._executar_geracao_e_pdf(
            gerar_proposta,
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        if pdf_prop: print("Proposta e PDF gerados.")
        else: print("Falha ao gerar Proposta ou seu PDF.")

        # 2. Gerar Material + PDF
        print("\n[2/2] Gerando Material Técnico...")
        _, pdf_mat = self._executar_geracao_e_pdf(
            gerar_material,
            abas_ativas,
            self.nome_closer_var.get(),
            self.celular_closer_var.get(),
            self.email_closer_var.get()
        )
        if pdf_mat: print("Material Técnico e PDF gerados.")
        else: print("Falha ao gerar Material Técnico ou seu PDF.")


        print("\n--- Geração de TUDO Concluída ---")
        # Mensagem final pode depender se ambos PDFs foram gerados com sucesso
        if pdf_prop and pdf_mat:
             showinfo("Concluído", "Proposta e Material Técnico (PPTX e PDF) gerados com sucesso.")
        else:
             showinfo("Concluído com Alertas", "Geração concluída, mas houve falha na criação de um ou mais arquivos PDF. Verifique o console para detalhes.")


# --- Função Principal (Execução) ---
def main():
    """Função principal que inicia a aplicação Tkinter."""
    print("Iniciando Gerador de Propostas...")
    app = MainApp()
    app.mainloop()

if __name__ == "__main__":
    # Adiciona tratamento de exceção global básico para capturar erros não tratados
    try:
        main()
    except Exception as e:
        print(f"\nERRO NÃO TRATADO NA APLICAÇÃO:\n{'-'*30}\n {type(e).__name__}: {e}\n{'-'*30}")
        # Tenta mostrar um erro para o usuário também
        try:
            root = tk.Tk()
            root.withdraw() # Esconde a janela root principal vazia
            showerror("Erro Crítico", f"Ocorreu um erro inesperado e a aplicação precisa fechar.\n\nDetalhes:\n{type(e).__name__}: {e}")
            root.destroy()
        except Exception:
            pass # Ignora se nem o Tkinter puder mostrar o erro
        sys.exit(1) # Termina com código de erro